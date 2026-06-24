"""Generate the investor-facing Xuannv business plan deck.

The deck is intentionally content-first: it uses a calm pastel theme, Chinese
copy, source-backed benchmark data, and project-local experiment visuals.
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "docs" / "presentation_assets"
OUT_DIR = ROOT / "docs" / "presentations"
OUT_PATH = OUT_DIR / "玄女科技BP_投资人版_v0.2.pptx"

COVER_IMAGE = ASSET_DIR / "geo_embedding_cover_pastel.png"
VIS_IMAGE = Path(
    "/data/xuannv_embedding/outputs/downstream/visualizations/"
    "harbin_stage2_v1_fold0/patch_000021_visualization.png"
)


class Theme:
    bg = RGBColor(250, 248, 241)
    ink = RGBColor(36, 51, 57)
    muted = RGBColor(102, 118, 125)
    green = RGBColor(137, 183, 164)
    blue = RGBColor(139, 180, 202)
    lavender = RGBColor(180, 170, 210)
    sand = RGBColor(235, 226, 204)
    line = RGBColor(216, 226, 220)
    white = RGBColor(255, 255, 255)
    dark_green = RGBColor(74, 121, 104)


FONT_CN = "Microsoft YaHei"
FONT_EN = "Aptos"


def ensure_dirs() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    OUT_DIR.mkdir(parents=True, exist_ok=True)


def crop_visual_panels() -> list[Path]:
    """Crop clean image panels from the existing English-labeled visualization."""

    if not VIS_IMAGE.exists():
        return []

    panel_specs = [
        ("s2_rgb", "S2 影像", (15, 120, 392, 497)),
        ("highres_optical", "高分光学", (459, 120, 836, 497)),
        ("s1_sar", "SAR / S1", (903, 120, 1280, 497)),
        ("landsat", "Landsat", (1347, 120, 1724, 497)),
        ("worldcover", "土地覆盖", (15, 528, 392, 905)),
        ("embedding_pca", "嵌入 PCA", (459, 528, 836, 905)),
        ("prediction", "预测概率", (903, 528, 1280, 905)),
    ]
    image = Image.open(VIS_IMAGE).convert("RGB")
    paths: list[Path] = []
    for slug, _label, box in panel_specs:
        out = ASSET_DIR / f"harbin_{slug}.png"
        image.crop(box).save(out)
        paths.append(out)
    return paths


def set_slide_bg(slide, color=Theme.bg) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 24,
    color: RGBColor = Theme.ink,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = FONT_CN,
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title: str, subtitle: str | None = None, kicker: str | None = None):
    if kicker:
        add_text(slide, kicker, 0.72, 0.42, 6.8, 0.24, 9, Theme.dark_green, True)
    add_text(slide, title, 0.72, 0.64, 10.9, 0.62, 25, Theme.ink, True)
    if subtitle:
        add_text(slide, subtitle, 0.74, 1.24, 11.0, 0.38, 10, Theme.muted)


def add_footer(slide, page: int, note: str = "玄女科技商业计划书 v0.2") -> None:
    add_text(slide, note, 0.72, 7.14, 6.0, 0.22, 7, Theme.muted)
    add_text(slide, f"{page:02d}", 12.0, 7.11, 0.55, 0.25, 8, Theme.muted, align=PP_ALIGN.RIGHT)


def add_round_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill=Theme.white,
    line=Theme.line,
    radius=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
):
    shape = slide.shapes.add_shape(radius, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    return shape


def add_chip(slide, text: str, x: float, y: float, w: float, color=Theme.green):
    shape = add_round_rect(slide, x, y, w, 0.34, RGBColor(240, 247, 244), color)
    add_text(slide, text, x + 0.08, y + 0.06, w - 0.16, 0.16, 8, Theme.dark_green, True, PP_ALIGN.CENTER)
    return shape


def add_card(slide, title: str, body: str, x: float, y: float, w: float, h: float, accent=Theme.green):
    add_round_rect(slide, x, y, w, h)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, title, x + 0.24, y + 0.18, w - 0.44, 0.33, 13, Theme.ink, True)
    add_text(slide, body, x + 0.24, y + 0.62, w - 0.44, h - 0.75, 9, Theme.muted, line_spacing=1.1)


def add_metric(slide, label: str, value: str, note: str, x: float, y: float, w: float, accent=Theme.blue):
    add_round_rect(slide, x, y, w, 1.04, RGBColor(253, 253, 250), Theme.line)
    add_text(slide, value, x + 0.18, y + 0.12, w - 0.32, 0.34, 19, accent, True)
    add_text(slide, label, x + 0.18, y + 0.52, w - 0.32, 0.22, 8, Theme.ink, True)
    add_text(slide, note, x + 0.18, y + 0.74, w - 0.32, 0.18, 7, Theme.muted)


def add_flow(slide, items: list[str], x: float, y: float, w: float, color=Theme.green) -> None:
    gap = 0.15
    bw = (w - gap * (len(items) - 1)) / len(items)
    for i, item in enumerate(items):
        bx = x + i * (bw + gap)
        add_round_rect(slide, bx, y, bw, 0.65, RGBColor(245, 250, 248), color)
        add_text(slide, item, bx + 0.06, y + 0.18, bw - 0.12, 0.20, 8, Theme.ink, True, PP_ALIGN.CENTER)
        if i < len(items) - 1:
            add_text(slide, "→", bx + bw + 0.02, y + 0.20, 0.12, 0.2, 10, Theme.muted, True)


def add_table_like(slide, rows: list[list[str]], x: float, y: float, col_ws: list[float], row_h: float = 0.56):
    total_w = sum(col_ws)
    for r, row in enumerate(rows):
        fill = RGBColor(241, 248, 246) if r == 0 else Theme.white
        add_round_rect(slide, x, y + r * row_h, total_w, row_h - 0.04, fill, Theme.line, MSO_AUTO_SHAPE_TYPE.RECTANGLE)
        cx = x
        for c, txt in enumerate(row):
            add_text(
                slide,
                txt,
                cx + 0.08,
                y + r * row_h + 0.12,
                col_ws[c] - 0.12,
                row_h - 0.22,
                7 if r else 8,
                Theme.ink if r == 0 else Theme.muted,
                bold=(r == 0 or c == 0),
            )
            cx += col_ws[c]


def create_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Cover
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    if COVER_IMAGE.exists():
        slide.shapes.add_picture(str(COVER_IMAGE), Inches(5.25), Inches(0.0), width=Inches(8.1), height=Inches(7.5))
    add_chip(slide, "商业计划书 · 投资人版", 0.75, 0.76, 1.82, Theme.green)
    add_text(slide, "玄女科技", 0.78, 1.55, 4.2, 0.56, 25, Theme.ink, True)
    add_text(slide, "中国地球空间智能底座", 0.78, 2.15, 5.6, 0.9, 30, Theme.dark_green, True)
    add_text(slide, "用地理嵌入赋能遥感智能应用。\n一次生成地球表征，多任务复用。", 0.82, 3.22, 4.9, 0.82, 15, Theme.muted)
    add_text(slide, "融资需求：5000 万元｜用于样板城市、模型工程化、任务库与商业交付", 0.82, 6.42, 6.9, 0.28, 9, Theme.muted)

    # 2 Scene pain
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "一个新区每月都要知道：哪里在施工、哪里占地、哪里发生变化", "传统遥感项目并不缺数据，真正缺的是把海量影像快速变成可核查变化线索的能力。", "01｜客户场景")
    add_metric(slide, "传统项目", "每类任务重跑一次", "找数据、预处理、标注、训练、核查、出报告", 0.9, 1.75, 3.5, Theme.blue)
    add_metric(slide, "玄女方案", "一个底座，多任务复用", "先生成区域地理嵌入，再调用施工、农用地、建筑等任务", 4.9, 1.75, 3.5, Theme.green)
    add_metric(slide, "客户价值", "先筛候选，再人工核查", "把人从全量盯图转到高优先级变化证据链", 8.9, 1.75, 3.5, Theme.lavender)
    add_text(slide, "核心问题", 0.95, 3.75, 1.5, 0.26, 10, Theme.muted, True)
    add_text(slide, "如何把海量地球观测数据压缩成可复用的地理嵌入，让变化检测、分类、预测和问答不再每个任务重做一遍？", 0.98, 4.16, 11.1, 0.8, 25, Theme.ink, True)
    add_flow(slide, ["海量观测", "地理嵌入", "多任务候选", "人工核查", "业务闭环"], 1.25, 5.9, 10.5, Theme.green)
    add_footer(slide, 2)

    # 3 Why now and market
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "现在是窗口期：商业航天让“看见”变便宜，地理智能决定“理解”效率", "以下数字代表广义产业基础、上游数据供给和地理空间分析大盘，不等同于玄女可直接获得市场。", "02｜为什么现在")
    add_metric(slide, "中国地理信息产业", "8501 亿元", "2024 年总产值，广义产业基础", 0.82, 1.68, 3.42, Theme.green)
    add_metric(slide, "地球观测小卫星市场", "26.4→55.2 亿美元", "2025-2030，上游数据供给市场", 4.85, 1.68, 3.42, Theme.blue)
    add_metric(slide, "全球地理空间分析", "1027→2340 亿美元", "2025-2033，广义分析市场", 8.88, 1.68, 3.42, Theme.lavender)
    add_card(slide, "美国对标给出的信号", "Planet、BlackSky、Maxar 证明遥感数据和实时监测存在持续需求；Esri 证明地理平台能进入政企核心工作流；Google 卫星嵌入证明地理表征层正在出现。", 0.82, 3.28, 5.55, 1.72, Theme.blue)
    add_card(slide, "中国需要自己的底座", "中国数据、政企流程、国产算力和安全合规要求不同；玄女借鉴的是地理表征层趋势，不复制美国公司的资产结构。", 6.82, 3.28, 5.55, 1.72, Theme.green)
    add_text(slide, "投资判断：未来竞争不只是谁拥有影像，而是谁能把影像持续变成可调用的地球智能。", 1.05, 5.82, 11.0, 0.42, 18, Theme.dark_green, True, PP_ALIGN.CENTER)
    add_footer(slide, 3)

    # 4 Industry bottleneck
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "行业不是没有需求，而是需求越多，传统交付方式越重", "今天遥感公司主要靠项目、人力和数据交付赚钱；换区域、换任务、换传感器，就要重做样本和模型。", "03｜行业痛点")
    add_flow(slide, ["找数据", "预处理", "标注", "训练", "核查", "交付"], 1.0, 1.8, 11.1, Theme.blue)
    add_card(slide, "收入方式", "数据销售、项目交付、遥感/GIS 平台授权、专业服务、少量订阅或接口。", 0.82, 3.0, 3.42, 1.45, Theme.green)
    add_card(slide, "天花板", "项目越多，人力越多；标注、质检、调参和交付吃掉毛利。", 4.85, 3.0, 3.42, 1.45, Theme.blue)
    add_card(slide, "真实证据", "天津遥感建筑半自动标注案例：AI 辅助让标注效率提升 80%，年均节省人力成本超 200 万元。", 8.88, 3.0, 3.42, 1.45, Theme.lavender)
    add_text(slide, "该案例证明遥感标注环节存在明确降本空间；玄女要验证的是把这种效率提升从单点任务扩展到多任务复用。", 1.04, 5.48, 11.0, 0.58, 18, Theme.dark_green, True, PP_ALIGN.CENTER)
    add_footer(slide, 4, "案例来源：国家数据局天津市测绘院遥感建筑半自动标注案例")

    # 5 Users
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "四类天使用户的痛点，本质上都是“每个任务重做一遍”", "玄女优先服务有预算、有数据、有反复任务的客户，而不是泛泛做遥感平台。", "04｜天使用户")
    add_card(slide, "B 端遥感项目公司", "不缺项目，缺可复用能力。多项目并行时，标注、质检、跨区域迁移和交付成本压低毛利。", 0.76, 1.72, 2.85, 1.72, Theme.green)
    add_card(slide, "政府/新区管委会", "不缺卫星图，缺能进入流程的变化证据链：图斑优先级、历史过程、冲突提示、外业核查闭环。", 3.83, 1.72, 2.85, 1.72, Theme.blue)
    add_card(slide, "准 C 端地块用户", "农场、园区、能源走廊、土地业主不想学遥感，只想知道地块有没有异常、风险在哪里。", 6.90, 1.72, 2.85, 1.72, Theme.lavender)
    add_card(slide, "学校科研端", "每个课题重复下载数据、做标签、训模型；真正需要的是可复用数据底座和快速任务验证。", 9.97, 1.72, 2.85, 1.72, Theme.green)
    add_text(slide, "共同需求：把遥感数据处理、标注、模型适配和变化候选生成变成一层统一能力，而不是每个专项项目从零开始。", 1.02, 4.46, 11.05, 0.66, 21, Theme.ink, True, PP_ALIGN.CENTER)
    add_flow(slide, ["采购触发", "试点区域", "多任务验证", "年度服务", "平台接口"], 1.15, 5.86, 10.8, Theme.green)
    add_footer(slide, 5)

    # 6 Harbin workflow/accounting
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "以哈尔滨新区为例：传统变化检测是一条多人协同的项目流水线", "区长/管委会要看施工工地、农用地、建筑变化等任务，下属需要协调数据、标注、模型和核查。", "05｜场景账本")
    add_flow(slide, ["接到任务", "协调人员", "下载多源数据", "标注样本", "训练/调参", "人工核查", "出报告"], 0.78, 1.82, 11.9, Theme.blue)
    add_card(slide, "时间消耗", "等数据、等标注、等训练、等质检；不同任务并行时排队更明显。", 0.88, 3.12, 3.55, 1.38, Theme.blue)
    add_card(slide, "空间消耗", "原始影像、中间栅格、标注版本、模型输出重复存储，很难在下一个任务里直接复用。", 4.88, 3.12, 3.55, 1.38, Theme.green)
    add_card(slide, "算力消耗", "每个任务独立训练或调参，重复跑编码器和预处理，GPU/NPU 时间沉淀不成公共资产。", 8.88, 3.12, 3.55, 1.38, Theme.lavender)
    add_text(slide, "玄女切入点：先把区域底座做好，再让每个任务只新增轻量适配和核查，不再从零开始。", 0.95, 5.55, 10.9, 0.48, 17, Theme.dark_green, True, PP_ALIGN.CENTER)
    add_footer(slide, 6)

    # 7 Solution
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "玄女方法：先为区域生成地理嵌入，再在同一底座上调用多任务", "一次生成地球表征，多任务复用；把重复的人力、存储和算力消耗前置成可沉淀的模型资产。", "06｜解决方案")
    add_flow(slide, ["多源观测", "统一地理嵌入", "任务头/少样本适配", "变化候选", "人工核查", "闭环沉淀"], 0.92, 1.76, 11.5, Theme.green)
    add_text(slide, "传统成本 = 每个任务 ×（数据 + 标注 + 训练 + 核查 + 交付）", 1.05, 3.15, 11.0, 0.36, 18, Theme.muted, True, PP_ALIGN.CENTER)
    add_text(slide, "玄女成本 = 区域底座一次生成 + 多任务轻量适配 + 重点核查", 1.05, 3.78, 11.0, 0.42, 22, Theme.dark_green, True, PP_ALIGN.CENTER)
    add_card(slide, "时间", "从“每个任务从零开始”变成“任务库调用 + 少样本适配”。", 0.88, 4.78, 2.82, 1.18, Theme.green)
    add_card(slide, "空间", "嵌入成为公共资产，减少重复中间成果和重复特征存储。", 3.92, 4.78, 2.82, 1.18, Theme.blue)
    add_card(slide, "算力", "编码器复用，任务头轻量化；后续重点在候选筛查和人机协同。", 6.96, 4.78, 2.82, 1.18, Theme.lavender)
    add_card(slide, "人力", "人不再盯全量影像，而是核查模型给出的高优先级变化。", 10.00, 4.78, 2.82, 1.18, Theme.green)
    add_footer(slide, 7)

    # 8 Token analogy
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "地理嵌入，是遥感行业进入智能应用前的“词元层”", "类比不是说遥感等于大语言模型，而是说：先有统一表征，才有可组合、可迁移、可调用的下游任务。", "07｜核心类比")
    add_card(slide, "大语言模型", "文本 → 词元 → 理解 / 检索 / 生成", 1.05, 2.0, 4.8, 1.55, Theme.blue)
    add_card(slide, "玄女", "地球观测 → 地理嵌入 → 变化检测 / 分类 / 预测 / 问答", 7.05, 2.0, 4.8, 1.55, Theme.green)
    add_text(slide, "GPT 先把文本转成词元，才能进行理解、检索、生成；玄女先把地球观测转成地理嵌入，才能进行变化检测、分类、预测和问答。", 1.18, 4.28, 10.9, 0.82, 22, Theme.ink, True, PP_ALIGN.CENTER)
    add_flow(slide, ["多源", "时序", "空间", "语义", "任务"], 2.0, 5.88, 9.3, Theme.green)
    add_footer(slide, 8, "参考趋势：Google Earth Engine 卫星嵌入 V1")

    # 9 Product
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "产品不是单点模型，而是地理嵌入底座 + 任务库 + 行业交付层", "先把区域数据资产化，再把任务从项目交付逐步产品化。", "08｜产品形态")
    add_card(slide, "底座层", "多源遥感数据接入、时空对齐、月度/区域地理嵌入生成、国产算力适配。", 1.0, 1.72, 3.2, 1.7, Theme.green)
    add_card(slide, "任务层", "施工工地、农用地变化、建筑变化、裸地/垃圾堆放、灾害扰动、水体变化等任务头。", 5.05, 1.72, 3.2, 1.7, Theme.blue)
    add_card(slide, "交付层", "图斑候选、证据链、接口、私有化部署、年度监测服务、人工核查闭环。", 9.1, 1.72, 3.2, 1.7, Theme.lavender)
    add_text(slide, "短期：样板城市 + 项目交付验证价值\n中期：任务库 + 私有化 / 年度服务提高复购\n长期：区域地理嵌入接口 / 订阅，进入政企数据栈", 1.15, 4.5, 10.8, 1.2, 19, Theme.ink, True, PP_ALIGN.CENTER)
    add_footer(slide, 9)

    # 10 Harbin sample
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "哈尔滨新区样板：同一区域底座，支撑多源输入与变化候选", "样板目标是形成可核查候选和成本账本，不把模型输出包装成自动执法结论。", "09｜样板验证")
    labels = ["S2 影像", "高分光学", "SAR/S1", "土地覆盖", "地理嵌入", "预测概率"]
    paths = [
        ASSET_DIR / "harbin_s2_rgb.png",
        ASSET_DIR / "harbin_highres_optical.png",
        ASSET_DIR / "harbin_s1_sar.png",
        ASSET_DIR / "harbin_worldcover.png",
        ASSET_DIR / "harbin_embedding_pca.png",
        ASSET_DIR / "harbin_prediction.png",
    ]
    for idx, (label, path) in enumerate(zip(labels, paths)):
        col = idx % 3
        row = idx // 3
        x, y = 0.88 + col * 2.05, 1.72 + row * 2.02
        add_text(slide, label, x, y - 0.20, 1.8, 0.18, 7, Theme.muted, True, PP_ALIGN.CENTER)
        if path.exists():
            slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(1.72), height=Inches(1.72))
    add_card(slide, "样板验证逻辑", "传统方式：每类任务重新组织数据、标注和模型。\n玄女方式：一次生成地理嵌入，施工、农用地、建筑变化等任务复用。\n交付定位：变化候选、优先级排序、证据链，而非替代最终判定。", 7.25, 1.72, 4.78, 3.75, Theme.green)
    add_text(slide, "定位：不是直接替代最终判定，而是让基层先看到最值得核查的变化。", 7.34, 5.78, 4.45, 0.42, 14, Theme.dark_green, True)
    add_footer(slide, 10)

    # 11 Competition
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "我们不和影像公司抢卫星，而是做任务复用的地理智能底座", "客户今天的替代方案各有价值，但很少同时解决多源时序、跨任务复用和私有化工作流。", "10｜竞争与替代")
    rows = [
        ["方案", "优势", "不足", "玄女定位"],
        ["传统项目制", "熟悉流程、可交付报告", "每个任务重做，人力重", "把重复劳动资产化"],
        ["影像/底图平台", "数据覆盖强", "偏数据供给，应用仍需二次开发", "连接上游数据与下游任务"],
        ["遥感/GIS 平台", "进入政企流程", "模型和任务复用不足", "补上 AI 表征和任务库"],
        ["通用遥感模型", "研发速度快", "本土数据、私有化、交付闭环不足", "面向中国场景工程化"],
    ]
    add_table_like(slide, rows, 0.72, 1.72, [1.6, 2.55, 3.1, 3.85], 0.72)
    add_text(slide, "壁垒来自三件事：区域级多源数据组织、可复用地理嵌入、进入客户核查流程的任务库。", 1.0, 5.95, 11.2, 0.4, 17, Theme.dark_green, True, PP_ALIGN.CENTER)
    add_footer(slide, 11)

    # 12 Business model and financials
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "商业模式：先用样板交付拿信任，再用任务库和年度服务提高复购", "以下为融资阶段的商业假设框架，单价和毛利需要随真实试点继续校准。", "11｜商业化")
    add_card(slide, "近期：样板项目", "城市/新区/园区样板，按区域、任务数量和交付周期收费；收入来自现有项目或平台升级预算。", 0.82, 1.75, 3.55, 1.38, Theme.green)
    add_card(slide, "中期：私有化 + 年度服务", "为遥感公司、政府部门和高校提供工具链、任务库、年度监测和技术服务。", 4.88, 1.75, 3.55, 1.38, Theme.blue)
    add_card(slide, "长期：接口/订阅", "区域地理嵌入、任务结果、样本检索和变化候选以接口或订阅进入政企数据栈。", 8.94, 1.75, 3.55, 1.38, Theme.lavender)
    add_metric(slide, "18 个月商业目标", "2-3 个样板", "形成可披露案例和成本账本", 0.9, 4.0, 3.35, Theme.green)
    add_metric(slide, "任务库目标", "5-10 类高频任务", "施工、农用地、建筑、灾害等", 4.95, 4.0, 3.35, Theme.blue)
    add_metric(slide, "收入模型", "项目 + 年服 + 接口", "从交付收入走向复购收入", 9.0, 4.0, 3.35, Theme.lavender)
    add_footer(slide, 12)

    # 13 Team/resources
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "团队与资源：同时理解遥感 AI、国产算力和政企交付", "核心团队将以遥感 AI、模型工程化、数据组织和政企交付资源形成复合能力。", "12｜团队")
    add_card(slide, "算法与工程", "多源时序遥感、地理嵌入、下游变化检测、模型工程化和国产 NPU 适配能力。", 0.9, 1.75, 3.55, 1.55, Theme.green)
    add_card(slide, "数据与场景", "已围绕哈尔滨、海淀等区域组织多源数据与研发验证，具备样板城市扩展基础。", 4.9, 1.75, 3.55, 1.55, Theme.blue)
    add_card(slide, "商业与交付", "目标客户明确：遥感项目公司、政府/新区管委会、科研机构和行业地块用户。", 8.9, 1.75, 3.55, 1.55, Theme.lavender)
    add_text(slide, "需要补强为投资人最终版的信息：创始人履历、团队人数、核心论文/专利/项目、客户资源和顾问背书。", 1.0, 4.7, 11.0, 0.62, 22, Theme.ink, True, PP_ALIGN.CENTER)
    add_footer(slide, 13)

    # 14 Financing
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "融资 5000 万：把地理嵌入从研发验证推到可复制样板", "资金重点投向样板闭环、模型工程化、任务库产品化和商业交付能力。", "13｜融资计划")
    add_metric(slide, "融资金额", "5000 万元", "轮次与估值待最终确认", 0.9, 1.72, 3.35, Theme.green)
    add_metric(slide, "周期", "18 个月", "从样板验证到可复制交付", 4.95, 1.72, 3.35, Theme.blue)
    add_metric(slide, "目标", "3 个样板 + 10 个任务", "用真实账本证明商业价值", 9.0, 1.72, 3.35, Theme.lavender)
    add_card(slide, "资金用途", "模型工程化与国产算力适配｜样板城市数据和标注｜产品化任务库｜销售与交付团队｜合规和数据采购。", 0.9, 3.35, 5.35, 1.6, Theme.green)
    add_card(slide, "18 个月里程碑", "拿下 2-3 个可披露样板；形成哈尔滨/海淀/雅江案例；沉淀 5-10 个高频任务；完成年度服务或私有化首单。", 6.95, 3.35, 5.35, 1.6, Theme.blue)
    add_text(slide, "最终要让投资人相信：这笔钱不是买更多实验，而是买“样板闭环 + 可复用资产 + 商业化入口”。", 1.02, 5.82, 11.1, 0.42, 16, Theme.ink, True, PP_ALIGN.CENTER)
    add_footer(slide, 14)

    # 15 Appendix capabilities
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "附录：研发验证显示已有候选排序信号，但仍需样板数据提升可交付精度", "对稀疏变化检测，AUC 说明排序信号；最佳 F1 和平均交并比仍低，当前定位为候选生成与优先级筛查。", "14｜附录 · 能力边界")
    cap_paths = [
        ("高分光学", ASSET_DIR / "harbin_highres_optical.png"),
        ("地理嵌入 PCA", ASSET_DIR / "harbin_embedding_pca.png"),
        ("预测概率", ASSET_DIR / "harbin_prediction.png"),
    ]
    for i, (label, path) in enumerate(cap_paths):
        x = 0.9 + i * 2.42
        add_text(slide, label, x, 1.82, 2.1, 0.2, 8, Theme.muted, True, PP_ALIGN.CENTER)
        if path.exists():
            slide.shapes.add_picture(str(path), Inches(x), Inches(2.08), width=Inches(2.05), height=Inches(2.05))
    add_card(slide, "研发信号", "哈尔滨研发集 5 折验证显示多个任务存在候选排序信号，AUC 约 0.75-0.89。", 8.45, 1.9, 3.75, 1.25, Theme.green)
    add_card(slide, "表达边界", "最佳 F1、固定阈值 F1 和平均交并比仍低；当前不作为自动判定或执法级结果，只做候选生成、人机核查和样板验证。", 8.45, 3.45, 3.75, 1.75, Theme.blue)
    add_footer(slide, 15)

    # 16 Appendix sources
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "附录：关键数据来源与口径", "用于投资人尽调追溯；正式路演可把详细链接放在资料包中。", "15｜附录 · 来源")
    add_card(slide, "海外对标", "Planet FY2026：收入 3.077 亿美元、经常性 ACV 98%\nBlackSky FY2025：收入 1.066 亿美元\nMaxar/Advent：交易估值约 64 亿美元\nEsri Fact Sheet：数十万组织、50% Fortune 500、多数国家政府", 0.82, 1.58, 5.65, 2.25, Theme.blue)
    add_card(slide, "市场与技术趋势", "Google Earth Engine 卫星嵌入 V1：全球 10 米、64 维、2017-2024\nMarketsandMarkets：EO 小卫星市场 2025-2030 CAGR 15.9%\nGrand View Research：地理空间分析市场 2025-2033 CAGR 10.4%\n国家数据局：天津遥感建筑半自动标注案例", 6.86, 1.58, 5.65, 2.25, Theme.green)
    add_text(slide, "完整链接：Planet、BlackSky、Maxar/Advent、Esri、Google Earth Engine、MarketsandMarkets、Grand View Research、国家数据局天津案例已记录在项目文档中。", 0.95, 4.58, 11.3, 0.52, 13, Theme.muted, align=PP_ALIGN.CENTER)
    add_text(slide, "口径提醒：广义产业和市场规模用于说明机会窗口，不等同于玄女可直接获得收入；玄女直接市场需结合样板客户、年度服务和接口订阅另行测算。", 1.0, 5.65, 11.1, 0.62, 17, Theme.dark_green, True, PP_ALIGN.CENTER)
    add_footer(slide, 16)

    return prs


def main() -> None:
    ensure_dirs()
    crop_visual_panels()
    prs = create_deck()
    prs.save(OUT_PATH)
    print(OUT_PATH)


if __name__ == "__main__":
    main()
