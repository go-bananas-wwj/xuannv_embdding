"""Generate the Xuannv capability case deck.

This follow-up deck starts after the first nine investor BP slides and proves
the foundation model through three concrete cases: Harbin New Area, Haidian,
and Yajiang hydropower monitoring.
"""

from __future__ import annotations

from pathlib import Path

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from create_front8_deck_v2 import C, bg, title, text, rect, line, arrow, picture, picture_crop, picture_fit, claim


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs" / "玄女科技BP_能力案例_v0.2.pptx"

CLOUD = Path("/root/workspace/bp_ppt素材汇总/cloud_examples")
MULTI = Path("/root/workspace/bp_ppt素材汇总/multi_task_examples")
MANUAL_REPORT = Path("/root/workspace/哈尔滨新区城市管理卫星遥感监测报告")
MODEL_REPORT = Path("/root/workspace/xuannv/reports/harbin_202512_202605/pdfs")
HARBIN_REPORT_ROOT = Path("/root/workspace/xuannv/reports/harbin_202512_202605")
HAIDIAN = Path("/root/workspace/report/assets")
YAJIANG = ROOT / "assets" / "user_materials" / "yajiang_case"
YAJIANG_AGENT = Path("/root/.codex/attachments/071916f0-9ca0-4d8c-83f1-bcd404e7ea3c/e178b904e06c08479cd215ca230ff245.png")
REGION_OUTLINES = {
    "harbin": ROOT / "assets" / "generated" / "practice_region_harbin_whitefit.png",
    "haidian": ROOT / "assets" / "generated" / "practice_region_haidian_whitefit.png",
    "yajiang": ROOT / "assets" / "generated" / "practice_region_yajiang_whitefit.png",
}
TMP = Path("/tmp/xuannv_capability_deck_assets")
RED = RGBColor(220, 38, 38)
PALE_RED = RGBColor(254, 242, 242)


def section_label(slide, value: str, x: float, y: float, color=C.blue) -> None:
    text(slide, value, x, y, 2.5, 0.16, 9, color, True)
    line(slide, x, y + 0.28, x + 1.34, y + 0.28, color, 1.0)


def small_label(slide, value: str, x: float, y: float, w: float, color=C.body) -> None:
    text(slide, value, x, y, w, 0.13, 7, color, True, PP_ALIGN.CENTER)


def outline_callout(slide, label_text: str, x: float, y: float, w: float, h: float, color=RED) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.background()
    box.line.color.rgb = color
    box.line.width = Pt(2.0)
    rect(slide, x, y - 0.22, 0.92, 0.22, PALE_RED, color)
    text(slide, label_text, x + 0.06, y - 0.155, 0.80, 0.08, 6, color, True, PP_ALIGN.CENTER)


def image_pill(slide, label_text: str, x: float, y: float, color=C.blue, fill=C.pale_blue) -> None:
    rect(slide, x, y, 0.70, 0.24, fill, color)
    text(slide, label_text, x + 0.07, y + 0.075, 0.56, 0.08, 6, color, True, PP_ALIGN.CENTER)


def bullet(slide, head: str, body: str, x: float, y: float, color=C.blue, fill=C.pale_blue, w: float = 3.66) -> None:
    rect(slide, x, y, w, 0.78, C.white, C.line)
    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.20), Inches(y + 0.26), Inches(0.26), Inches(0.26))
    dot.fill.solid()
    dot.fill.fore_color.rgb = fill
    dot.line.color.rgb = color
    dot.line.width = Pt(1.0)
    text(slide, head, x + 0.62, y + 0.14, w - 0.92, 0.18, 10, C.ink, True)
    text(slide, body, x + 0.62, y + 0.44, w - 0.92, 0.18, 8, C.body)


def compact_bullet(slide, head: str, body: str, x: float, y: float, color=C.blue, fill=C.pale_blue, w: float = 3.30) -> None:
    rect(slide, x, y, w, 0.58, C.white, C.line)
    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.18), Inches(y + 0.19), Inches(0.20), Inches(0.20))
    dot.fill.solid()
    dot.fill.fore_color.rgb = fill
    dot.line.color.rgb = color
    dot.line.width = Pt(0.9)
    text(slide, head, x + 0.50, y + 0.10, w - 0.72, 0.14, 9, C.ink, True)
    text(slide, body, x + 0.50, y + 0.34, w - 0.72, 0.12, 7, C.body)


def flow_step(slide, index: int, value: str, x: float, y: float, color=C.blue, fill=C.pale_blue, w: float = 3.34) -> None:
    rect(slide, x, y, w, 0.34, C.white, C.line)
    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.10), Inches(y + 0.08), Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = fill
    dot.line.color.rgb = color
    dot.line.width = Pt(0.8)
    text(slide, f"{index}", x + 0.11, y + 0.105, 0.16, 0.06, 5, color, True, PP_ALIGN.CENTER)
    text(slide, value, x + 0.36, y + 0.09, w - 0.48, 0.12, 8, C.ink, True)


def flow_node(
    slide,
    value: str,
    x: float,
    y: float,
    color=C.blue,
    fill=C.pale_blue,
    w: float = 1.02,
    h: float = 0.62,
    duration: str = "",
    note: str = "",
) -> None:
    rect(slide, x, y, w, h, fill, color)
    text(slide, value, x + 0.06, y + 0.08, w - 0.12, 0.12, 7, color, True, PP_ALIGN.CENTER)
    if duration:
        text(slide, duration, x + 0.06, y + 0.27, w - 0.12, 0.10, 7, C.ink, True, PP_ALIGN.CENTER)
    if note:
        text(slide, note, x + 0.06, y + 0.44, w - 0.12, 0.08, 5, C.muted, True, PP_ALIGN.CENTER)


def metric(slide, value: str, label: str, x: float, y: float, color=C.blue, fill=C.pale_blue, w: float = 1.72) -> None:
    rect(slide, x, y, w, 0.78, fill, color)
    text(slide, value, x + 0.10, y + 0.13, w - 0.20, 0.22, 15, color, True, PP_ALIGN.CENTER)
    text(slide, label, x + 0.10, y + 0.47, w - 0.20, 0.14, 7, C.body, True, PP_ALIGN.CENTER)


def case_card(slide, name: str, scene: str, need: str, image: Path, x: float, color, fill) -> None:
    text(slide, name, x + 0.20, 1.90, 3.25, 0.28, 17, C.ink, True, PP_ALIGN.CENTER)
    picture_fit(slide, image, x + 0.06, 2.26, 3.56, 2.42)
    line(slide, x + 0.32, 4.96, x + 3.36, 4.96, color, 1.1)
    rect(slide, x + 0.26, 5.22, 3.14, 0.48, fill, color)
    text(slide, scene, x + 0.44, 5.37, 2.78, 0.14, 10, color, True, PP_ALIGN.CENTER)
    text(slide, need, x + 0.26, 5.94, 3.14, 0.40, 12, C.ink, True, PP_ALIGN.CENTER)


def task_thumb(slide, label: str, path: Path, x: float, y: float, w: float, h: float) -> None:
    picture_fit(slide, path, x, y, w, h)
    small_label(slide, label, x, y + h + 0.08, w)


def render_pdf_page(pdf: Path, page: int, name: str, zoom: float = 1.55) -> Path:
    TMP.mkdir(parents=True, exist_ok=True)
    out = TMP / f"{name}_p{page + 1}.png"
    if out.exists():
        return out
    doc = fitz.open(pdf)
    pix = doc[page].get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
    pix.save(out)
    return out


def report_page_counts(folder: Path) -> int:
    total = 0
    for pdf in folder.glob("*.pdf"):
        with fitz.open(pdf) as doc:
            total += doc.page_count
    return total


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    cloud_imgs = sorted(CLOUD.glob("*.png"))
    multi_imgs = {
        "水体识别": MULTI / "patch_000041_jrc_water_2025-04_detail.png",
        "建筑提取": MULTI / "patch_000354_building_extraction_2025-04_detail.png",
        "施工工地": MULTI / "patch_000217_construction_2025-08_vs_2025-09_detail.png",
        "耕地变化": MULTI / "patch_000217_land_conversion_2025-08_vs_2025-09_detail.png",
        "地表分类": MULTI / "patch_000339_dynamic_world_2025-04_detail.png",
    }
    manual_total = MANUAL_REPORT / "00-2025年下半年哈尔滨新区城市管理卫星遥感监测总报告.pdf"
    model_total = MODEL_REPORT / "00-2025年12月-2026年5月哈尔滨新区城市管理卫星遥感监测总报告.pdf"
    manual_cover = render_pdf_page(manual_total, 0, "manual_total")
    model_cover = render_pdf_page(model_total, 0, "model_total")
    model_construction_page = render_pdf_page(
        MODEL_REPORT / "01-2025年12月-2026年5月哈尔滨新区建筑工地卫星遥感监测报告.pdf", 5, "model_construction_p6"
    )
    model_building_page = render_pdf_page(
        MODEL_REPORT / "02-2025年12月-2026年5月哈尔滨新区建筑变化卫星遥感监测报告.pdf", 5, "model_building_p6"
    )
    model_farmland_page = render_pdf_page(
        MODEL_REPORT / "03-2025年12月-2026年5月哈尔滨新区耕地非农非粮卫星遥感监测报告.pdf", 5, "model_farmland_p6"
    )
    model_trash_page = render_pdf_page(
        MODEL_REPORT / "04-2025年12月-2026年5月哈尔滨新区无序堆放垃圾渣土卫星遥感监测报告.pdf", 5, "model_trash_p6"
    )
    manual_pages = report_page_counts(MANUAL_REPORT)
    model_pages = report_page_counts(MODEL_REPORT)
    harbin_map = HARBIN_REPORT_ROOT / "pdf_maps" / "all_categories_map.png"
    haidian_chart = HAIDIAN / "aef_vs_v1.0_v2.png"
    haidian_patch = HAIDIAN / "compare_haidian_construction_patch_000198.png"

    # 10. Overview
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "10", "现有工程实践案例")
    text(s, "玄女底座已围绕政府城市治理、复杂城市变化评测与水电站周边监测，形成面向真实区域的工程实践。", 1.02, 1.18, 11.20, 0.34, 15, C.blue, True, PP_ALIGN.CENTER)
    case_card(s, "哈尔滨新区", "政府客户 · 城市治理", "多专题持续监测与城市管理报告更新", REGION_OUTLINES["harbin"], 0.86, C.blue, C.pale_blue)
    case_card(s, "海淀区", "城市核心区 · 任务评测", "复杂城市纹理下，对标国际地理嵌入能力", REGION_OUTLINES["haidian"], 4.82, C.green, C.mint)
    case_card(s, "雅江区域", "水电站周边 · 安全监测", "水电站周边态势监测与遥感分析报告", REGION_OUTLINES["yajiang"], 8.78, C.purple, C.pale_purple)

    # 11. Harbin cloud robustness
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "11", "哈尔滨新区：有云遮挡，仍能识别变化区域")
    text(s, "政府客户的真实难题", 0.62, 1.24, 3.08, 0.26, 17, C.blue, True)
    line(s, 0.62, 1.66, 2.42, 1.66, C.blue, 1.1)
    text(s, "云、雾、雪遮挡关键地物，城市治理仍需要连续、稳定的变化检测结果。", 0.62, 1.90, 3.12, 0.40, 14, C.ink, True)
    compact_bullet(s, "缺失模态下仍可推理", "不依赖单张完美影像完成判断", 0.62, 2.58, C.blue, C.pale_blue, 3.12)
    compact_bullet(s, "跨模态保留变化信号", "RGB 与嵌入共同支撑变化检测", 0.62, 3.25, C.green, C.mint, 3.12)
    compact_bullet(s, "减少人工重跑", "影像质量波动不直接造成流程停摆", 0.62, 3.92, C.purple, C.pale_purple, 3.12)
    rect(s, 0.62, 4.78, 3.12, 1.28, C.pale_blue, C.blue)
    text(s, "为什么能识别？", 0.82, 4.96, 2.72, 0.14, 10, C.blue, True, PP_ALIGN.CENTER)
    text(s, "训练阶段融合多模态遥感数据；单一模态被云雾遮挡时，其他模态仍可补充约束。对比变化前后地理嵌入差异，即可识别真实变化区域。", 0.82, 5.23, 2.72, 0.42, 8, C.body, True, PP_ALIGN.CENTER)
    line(s, 4.16, 1.40, 4.16, 6.46, C.line, 0.9)
    col_labels = ["变化前 RGB", "变化后 RGB", "变化前嵌入", "变化后嵌入", "变化概率"]
    for j, label_value in enumerate(col_labels):
        text(s, label_value, 5.18 + j * 1.42, 1.45, 1.34, 0.12, 7, C.muted, True, PP_ALIGN.CENTER)
    sample_names = ["样例一", "样例二", "样例三"]
    for i, img in enumerate(cloud_imgs):
        y = 1.70 + i * 1.60
        picture(s, img, 5.18, y, 7.10, 1.62)
        image_pill(s, sample_names[i], 4.34, y + 0.68, C.blue if i == 0 else (C.green if i == 1 else C.purple), C.pale_blue if i == 0 else (C.mint if i == 1 else C.pale_purple))
    outline_callout(s, "云雾遮挡", 6.76, 1.86, 1.18, 0.78, RED)
    outline_callout(s, "云雾遮挡", 5.28, 3.36, 1.32, 0.90, RED)
    outline_callout(s, "云雾遮挡", 5.28, 5.10, 1.30, 0.78, RED)
    claim(s, "玄女底座提高的是业务可用性：真实遥感数据不完美，但系统仍要稳定产出。", 7.02, C.blue)

    # 12. Harbin multi-task reuse
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "12", "哈尔滨新区：同一套地理嵌入，复用到多类城市治理任务")
    text(s, "政府客户的真实难题", 0.62, 1.24, 3.08, 0.26, 17, C.blue, True)
    line(s, 0.62, 1.66, 2.42, 1.66, C.blue, 1.1)
    text(s, "城市治理不是单一模型问题，而是水体、建筑、地表、施工工地和耕地变化等多任务并行。", 0.62, 1.90, 3.12, 0.46, 14, C.ink, True)
    compact_bullet(s, "旧标签复用", "旧标签局部过时，也能提供可迁移监督", 0.62, 2.64, C.blue, C.pale_blue, 3.12)
    compact_bullet(s, "一次表征，多次复用", "同一区域嵌入可服务多类任务头", 0.62, 3.31, C.green, C.mint, 3.12)
    compact_bullet(s, "支持多类下游任务", "从水体、建筑扩展到变化与分类任务", 0.62, 3.98, C.purple, C.pale_purple, 3.12)
    rect(s, 0.62, 4.84, 3.12, 1.10, C.pale_blue, C.blue)
    text(s, "核心价值", 0.82, 5.02, 2.72, 0.14, 10, C.blue, True, PP_ALIGN.CENTER)
    text(s, "从旧标签中学习稳定地物知识，在新影像上完成更贴近真实地物的划分。", 0.82, 5.27, 2.72, 0.34, 8, C.body, True, PP_ALIGN.CENTER)
    line(s, 4.16, 1.32, 4.16, 5.94, C.line, 0.9)
    text(s, "旧标签复用后的新影像划分效果", 4.62, 1.18, 7.58, 0.18, 13, C.blue, True, PP_ALIGN.CENTER)
    positions = [
        ("水体识别", multi_imgs["水体识别"], 4.50, 1.48, 3.82, 1.06),
        ("建筑提取", multi_imgs["建筑提取"], 4.50, 2.98, 3.82, 1.06),
        ("地表分类", multi_imgs["地表分类"], 4.50, 4.48, 3.82, 1.06),
        ("施工工地变化", multi_imgs["施工工地"], 8.58, 1.68, 3.82, 1.48),
        ("耕地非农非粮", multi_imgs["耕地变化"], 8.58, 3.88, 3.82, 1.48),
    ]
    for name, path, x, y, w, h in positions:
        task_thumb(s, name, path, x, y, w, h)
    text(s, "旧标签不再只是一次性标注资产：玄女底座从旧标签中提取可迁移地物知识，并将同一套地理嵌入复用到水体、建筑、地表分类与变化检测等多类任务。", 0.86, 6.24, 11.60, 0.24, 12, C.blue, True, PP_ALIGN.CENTER)

    # 13. Harbin report delivery
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "13", "哈尔滨新区：从检测结果到政府可用报告")
    text(s, "模型价值最终要落到交付：把变化检测、统计制图和业务解释组织成政府客户可阅读、可复核、可归档的报告。", 1.02, 1.14, 11.20, 0.30, 13, C.body, False, PP_ALIGN.CENTER)
    text(s, "用户旅程对比", 0.62, 1.62, 3.60, 0.20, 16, C.blue, True)
    line(s, 0.62, 1.98, 2.20, 1.98, C.blue, 1.0)
    text(s, "传统人工流程", 0.62, 2.12, 1.56, 0.16, 10, C.ink, True)
    text(s, "10-15 天 / 多人串行", 2.22, 2.12, 1.54, 0.16, 8, C.amber, True, PP_ALIGN.RIGHT)
    manual_nodes = [
        ("接到任务", "0.5 天", "拆需求", 0.54, 2.42),
        ("下载影像", "1-2 天", "等数据", 1.72, 2.42),
        ("预处理", "1-2 天", "配准清洗", 2.90, 2.42),
        ("人工解译", "5-8 天", "逐图核查", 2.90, 3.30),
        ("制图统计", "2-3 天", "反复整理", 1.72, 3.30),
        ("审核归档", "1-2 天", "多轮返工", 0.54, 3.30),
    ]
    for value, duration, note, x, y in manual_nodes:
        flow_node(s, value, x, y, C.amber, C.pale_amber, duration=duration, note=note)
    arrow(s, 1.59, 2.73, 1.69, 2.73, C.amber, 1.35)
    arrow(s, 2.77, 2.73, 2.87, 2.73, C.amber, 1.35)
    arrow(s, 3.41, 3.07, 3.41, 3.27, C.amber, 1.35)
    arrow(s, 2.87, 3.61, 2.77, 3.61, C.amber, 1.35)
    arrow(s, 1.69, 3.61, 1.59, 3.61, C.amber, 1.35)
    text(s, "玄女自动化流程", 0.62, 4.32, 1.56, 0.16, 10, C.blue, True)
    text(s, "0.5-1 天 / 底座复用", 2.18, 4.32, 1.58, 0.16, 8, C.blue, True, PP_ALIGN.RIGHT)
    xuannv_nodes = [
        ("选择任务", "5 分钟", "低门槛", 0.54, 4.62),
        ("调用底座", "自动推理", "免重训", 1.72, 4.62),
        ("自动统计", "分钟级", "自动制图", 2.90, 4.62),
        ("生成报告", "10 分钟", "自动成文", 1.72, 5.48),
        ("人工复核", "0.5 天", "确认提交", 0.54, 5.48),
    ]
    for value, duration, note, x, y in xuannv_nodes:
        flow_node(s, value, x, y, C.blue, C.pale_blue, duration=duration, note=note)
    arrow(s, 1.59, 4.93, 1.69, 4.93, C.blue, 1.35)
    arrow(s, 2.77, 4.93, 2.87, 4.93, C.blue, 1.35)
    line(s, 3.41, 5.27, 3.41, 5.79, C.blue, 1.35)
    arrow(s, 3.41, 5.79, 2.77, 5.79, C.blue, 1.35)
    text(s, "↓", 3.31, 5.36, 0.20, 0.12, 10, C.blue, True, PP_ALIGN.CENTER)
    text(s, "←", 3.00, 5.62, 0.20, 0.12, 10, C.blue, True, PP_ALIGN.CENTER)
    arrow(s, 1.69, 5.79, 1.59, 5.79, C.blue, 1.35)
    line(s, 4.18, 1.48, 4.18, 6.58, C.line, 0.9)
    text(s, "自动生成报告样例", 4.42, 1.62, 4.10, 0.20, 14, C.blue, True, PP_ALIGN.CENTER)
    report_tiles = [
        ("总报告", model_cover, 4.42, 2.02),
        ("建筑工地", model_construction_page, 5.78, 2.02),
        ("建筑变化", model_building_page, 7.14, 2.02),
        ("耕地非农非粮", model_farmland_page, 4.42, 3.98),
        ("垃圾渣土", model_trash_page, 5.78, 3.98),
        ("点位分布", harbin_map, 7.14, 3.98),
    ]
    for label, img, x, y in report_tiles:
        picture_fit(s, img, x, y, 1.10, 1.48)
        small_label(s, label, x, y + 1.60, 1.10, C.blue)
    text(s, f"已生成 5 份专题报告，共 {model_pages} 页，覆盖多类城市治理任务。", 4.42, 5.98, 4.10, 0.24, 8, C.body, True, PP_ALIGN.CENTER)
    line(s, 8.70, 1.48, 8.70, 6.58, C.line, 0.9)
    text(s, "交付效率测算", 9.10, 1.62, 3.20, 0.22, 16, C.blue, True, PP_ALIGN.CENTER)
    text(s, "以一次月度/周期监测报告为例，按遥感工程师、制图与报告撰写人力的保守估算。", 9.04, 2.02, 3.28, 0.34, 9, C.body, True, PP_ALIGN.CENTER)
    metric(s, "10-15 天", "传统交付周期", 9.02, 2.66, C.amber, C.pale_amber, 1.50)
    metric(s, "4-6 万元", "传统综合成本", 10.78, 2.66, C.amber, C.pale_amber, 1.50)
    metric(s, "0.5-1 天", "玄女交付周期", 9.02, 3.60, C.blue, C.pale_blue, 1.50)
    metric(s, "千元级", "玄女边际成本", 10.78, 3.60, C.blue, C.pale_blue, 1.50)
    rect(s, 9.02, 4.62, 3.26, 0.92, C.pale_blue, C.blue)
    text(s, "节省结果", 9.20, 4.78, 2.90, 0.14, 10, C.blue, True, PP_ALIGN.CENTER)
    text(s, "报告生产从“项目制人力交付”变成“模型驱动的持续交付”。", 9.22, 5.06, 2.86, 0.22, 8, C.body, True, PP_ALIGN.CENTER)
    text(s, "时间节省约 90%+，单次报告成本下降一个数量级。", 0.86, 6.94, 11.60, 0.22, 14, C.blue, True, PP_ALIGN.CENTER)

    # 14. Haidian benchmark
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "14", "海淀区 Benchmark：多时序变化任务，玄女 V1.0 优于 AEF")
    text(s, "用真实下游任务指标说话：海淀区 5-fold CV，对比玄女 V1.0 embedding 与 AEF 2025 官方年度 embedding。", 1.02, 1.14, 11.20, 0.30, 13, C.body, False, PP_ALIGN.CENTER)
    rect(s, 0.82, 1.72, 3.32, 1.12, C.pale_blue, C.blue)
    text(s, "核心结论", 1.06, 1.92, 2.84, 0.16, 11, C.blue, True, PP_ALIGN.CENTER)
    text(s, "在建筑变化、耕地变化、垃圾堆放等多时序变化类任务上，玄女 V1.0 表现更优。", 1.02, 2.22, 2.92, 0.28, 9, C.ink, True, PP_ALIGN.CENTER)
    rect(s, 0.82, 3.16, 3.32, 1.04, C.pale_amber, C.amber)
    text(s, "坦诚短板", 1.06, 3.34, 2.84, 0.16, 11, C.amber, True, PP_ALIGN.CENTER)
    text(s, "施工工地属于单时序识别类任务，当前训练数据不足，AEF 暂时更强。", 1.02, 3.64, 2.92, 0.28, 9, C.ink, True, PP_ALIGN.CENTER)
    rect(s, 0.82, 4.52, 3.32, 0.82, C.white, C.line)
    text(s, "怎么看这张图？", 1.06, 4.68, 2.84, 0.14, 10, C.blue, True, PP_ALIGN.CENTER)
    text(s, "蓝色为玄女 V1.0，红色为 AEF；同一任务下柱越高越好。", 1.02, 4.96, 2.92, 0.18, 8, C.body, True, PP_ALIGN.CENTER)
    line(s, 4.42, 1.56, 4.42, 5.78, C.line, 0.9)
    text(s, "四类任务、四组指标对比", 4.72, 1.58, 7.24, 0.18, 13, C.blue, True, PP_ALIGN.CENTER)
    picture_fit(s, haidian_chart, 4.58, 1.96, 7.62, 2.72)
    text(s, "变化类任务：建筑变化 / 耕地变化 / 垃圾堆放", 4.80, 4.92, 3.84, 0.16, 9, C.blue, True, PP_ALIGN.CENTER)
    text(s, "单时序识别任务：施工工地", 9.04, 4.92, 2.76, 0.16, 9, C.amber, True, PP_ALIGN.CENTER)
    metric(s, "0.8828", "建筑变化 AUC", 4.70, 5.48, C.blue, C.pale_blue, 1.46)
    metric(s, "0.8876", "耕地变化 AUC", 6.46, 5.48, C.blue, C.pale_blue, 1.46)
    metric(s, "0.8876", "垃圾堆放 AUC", 8.22, 5.48, C.blue, C.pale_blue, 1.46)
    metric(s, "AEF 更强", "施工工地任务", 9.98, 5.48, C.amber, C.pale_amber, 1.46)
    text(s, "注：5-fold CV，海淀区下游任务评测；指标包括 AUC-ROC、F1-best、F1@0.5、mIoU。", 1.00, 6.42, 11.30, 0.16, 8, C.muted, False, PP_ALIGN.CENTER)
    claim(s, "玄女已经在中国城市变化类任务上形成可量化优势；单时序短板，可通过本土训练数据继续补齐。", 6.76, C.blue)

    # 15. Haidian business meaning
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "15", "海淀区：本土任务需要本土遥感嵌入底座")
    text(s, "海淀 Benchmark 的价值不只是跑分，而是说明：通用地理嵌入进入真实城市治理任务后，会受到本土数据、标签口径和业务目标的共同影响。", 1.02, 1.14, 11.20, 0.30, 13, C.body, False, PP_ALIGN.CENTER)
    text(s, "真实任务样例", 0.82, 1.62, 5.32, 0.20, 14, C.blue, True, PP_ALIGN.CENTER)
    picture_fit(s, haidian_patch, 0.86, 2.00, 5.20, 2.66)
    small_label(s, "高分影像 / 嵌入可视化 / 预测结果 / 标注对比", 0.86, 4.82, 5.20, C.blue)
    rect(s, 0.86, 5.34, 5.20, 0.56, C.pale_blue, C.blue)
    text(s, "这类城市建设活动具有明显的本土业务定义：什么算施工工地、变化是否有效、是否需要进入报告，都依赖本地标签体系。", 1.04, 5.50, 4.84, 0.20, 8, C.body, True, PP_ALIGN.CENTER)
    line(s, 6.42, 1.54, 6.42, 6.16, C.line, 0.9)
    text(s, "Benchmark 传递的三层信息", 6.86, 1.62, 5.10, 0.20, 14, C.blue, True, PP_ALIGN.CENTER)
    rect(s, 6.92, 2.08, 4.86, 0.76, C.pale_blue, C.blue)
    text(s, "1 方向已被验证", 7.12, 2.22, 4.46, 0.14, 10, C.blue, True, PP_ALIGN.CENTER)
    text(s, "AlphaEarth/AEF 证明地理嵌入正在成为遥感智能基础设施。", 7.14, 2.48, 4.42, 0.14, 8, C.body, True, PP_ALIGN.CENTER)
    arrow(s, 9.35, 2.90, 9.35, 3.20, C.line, 1.0)
    rect(s, 6.92, 3.28, 4.86, 0.86, C.mint, C.green)
    text(s, "2 本土任务产生差异", 7.12, 3.42, 4.46, 0.14, 10, C.green, True, PP_ALIGN.CENTER)
    text(s, "建筑变化、耕地变化、垃圾堆放等多时序任务，玄女已形成可量化优势。", 7.14, 3.68, 4.42, 0.18, 8, C.body, True, PP_ALIGN.CENTER)
    arrow(s, 9.35, 4.20, 9.35, 4.50, C.line, 1.0)
    rect(s, 6.92, 4.58, 4.86, 0.86, C.pale_amber, C.amber)
    text(s, "3 数据补齐形成壁垒", 7.12, 4.72, 4.46, 0.14, 10, C.amber, True, PP_ALIGN.CENTER)
    text(s, "施工工地等单时序任务暂时落后，反而明确了补数据、补标签、补任务定义的迭代路径。", 7.14, 4.98, 4.42, 0.18, 8, C.body, True, PP_ALIGN.CENTER)
    claim(s, "商业含义：不是简单国产替代，而是用本土数据和业务闭环，持续训练出更适配中国城市治理的地理智能底座。", 6.64, C.blue)

    # 16. Yajiang agent
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "16", "雅江区域：把地理嵌入接入水电站监测智能体")
    text(s, "雅江案例面向水电站监测：用户选择区域、任务和时间，用自然语言触发模型与下游任务，自动生成结构化遥感分析报告。", 1.02, 1.14, 11.20, 0.30, 13, C.body, False, PP_ALIGN.CENTER)
    picture_fit(s, YAJIANG_AGENT, 0.76, 1.70, 8.12, 4.58)
    bullet(s, "自然语言发起任务", "用户直接提出报告需求，而不是操作脚本", 9.20, 1.78, C.blue, C.pale_blue, 3.10)
    bullet(s, "模型与指标可追溯", "展示模型版本、精度、置信度、低置信比例", 9.20, 2.82, C.green, C.mint, 3.10)
    bullet(s, "报告可继续对话", "历史会话保留，可追加月份和修改任务", 9.20, 3.86, C.purple, C.pale_purple, 3.10)
    metric(s, "96.8%", "总体精度", 9.20, 5.12, C.blue, C.pale_blue, 1.40)
    metric(s, "96.3%", "平均置信度", 10.88, 5.12, C.green, C.mint, 1.40)
    claim(s, "玄女底座可以变成用户可直接使用的智能体产品界面，而不是停留在模型 API。", 6.58, C.blue)

    # 17. Yajiang downstream to report
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "17", "从下游任务结果，到可阅读的遥感分析报告")
    text(s, "雅江水电站监测需要的不只是预测图，而是能持续生成指标、图表、风险提示和后续行动建议的工作流。", 1.02, 1.14, 11.20, 0.30, 13, C.body, False, PP_ALIGN.CENTER)
    y_tasks = [
        ("雅江区域", YAJIANG / "yajiang_area.png", 0.86, 1.74, 1.82, 1.64),
        ("地物分类", YAJIANG / "land_classification.png", 3.10, 1.74, 2.48, 1.64),
        ("空间聚类", YAJIANG / "spatial_cluster.png", 6.00, 1.74, 2.48, 1.64),
        ("高程回归", YAJIANG / "elevation_regression.png", 8.90, 1.74, 2.48, 1.64),
        ("变化检测", YAJIANG / "change_detection.png", 1.32, 4.10, 1.62, 1.56),
        ("嵌入检索", YAJIANG / "embedding_retrieval.png", 3.64, 4.10, 1.62, 1.56),
        ("坡度风险", YAJIANG / "slope_risk.png", 5.96, 4.10, 1.62, 1.56),
    ]
    for name, path, x, y, w, h in y_tasks:
        task_thumb(s, name, path, x, y, w, h)
    rect(s, 8.70, 4.18, 3.26, 1.42, C.pale_blue, C.blue)
    text(s, "智能体报告", 9.00, 4.48, 2.66, 0.22, 14, C.blue, True, PP_ALIGN.CENTER)
    text(s, "汇总预测图、置信度、误差图、关键指标与行动建议", 9.02, 4.92, 2.62, 0.30, 9, C.body, True, PP_ALIGN.CENTER)
    arrow(s, 7.76, 4.86, 8.52, 4.86, C.line, 1.2)
    claim(s, "玄女的最终产品形态：地理嵌入底座 + 下游任务库 + 报告智能体。", 6.58, C.blue)

    # 18. Product capability matrix
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "18", "从三个案例抽象出玄女底座的四层产品能力")
    text(s, "案例不是孤立项目，而是在验证玄女可复制的产品架构。", 1.02, 1.14, 11.20, 0.30, 13, C.body, False, PP_ALIGN.CENTER)
    layers = [
        ("地理嵌入层", "多源、多时相、跨模态、缺失模态下可用", "哈尔滨云遮挡样例", C.blue, C.pale_blue),
        ("任务适配层", "分类、检测、变化监测、检索、风险识别共享底座", "哈尔滨多任务 / 雅江多任务", C.green, C.mint),
        ("评测对标层", "与 AEF 官方嵌入进行下游任务量化比较", "海淀 V1.0 Benchmark", C.amber, C.pale_amber),
        ("报告交付层", "生成政务报告或智能体报告，进入客户工作流", "哈尔滨 2026 报告 / 雅江智能体", C.purple, C.pale_purple),
    ]
    for i, (head, body, evidence, color, fill) in enumerate(layers):
        y = 1.78 + i * 1.02
        rect(s, 1.02, y, 11.30, 0.76, C.white, C.line)
        rect(s, 1.22, y + 0.18, 0.40, 0.40, fill, color)
        text(s, head, 1.88, y + 0.17, 1.74, 0.18, 12, color, True, PP_ALIGN.CENTER)
        text(s, body, 4.02, y + 0.17, 4.22, 0.20, 10, C.ink, True, PP_ALIGN.CENTER)
        text(s, evidence, 8.72, y + 0.17, 2.94, 0.20, 9, C.body, True, PP_ALIGN.CENTER)
    line(s, 1.72, 5.98, 11.64, 5.98, C.line, 0.8)
    text(s, "商业化承接", 1.62, 6.28, 2.00, 0.20, 12, C.ink, True, PP_ALIGN.CENTER)
    text(s, "按区域、按任务、按更新频率、按 API / 平台授权收费。", 3.92, 6.25, 6.70, 0.24, 12, C.blue, True, PP_ALIGN.CENTER)
    claim(s, "玄女不是遥感项目公司，而是在搭建可复用、可评测、可交付的地理智能基础设施。", 6.86, C.blue)

    return prs


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
