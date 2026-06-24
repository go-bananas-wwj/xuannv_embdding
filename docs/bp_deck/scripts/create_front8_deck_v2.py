"""Generate the revised first nine Xuannv BP slides.

This deck is a stage-presentation draft, not a document handout.  It avoids
internal notes, keeps each slide to one claim, and favors real project images
plus clean editable diagrams.
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
USER = ASSETS / "user_materials"
PPT_MEDIA = USER / "harbin_yajiang_ppt_media"
YAJIANG = USER / "yajiang_case"
AEF = ASSETS / "aef_references"
E4_AEF = USER / "e4_alphaearth_media"
OPENGEOSCOPE = ASSETS / "opengoscope"
CACHE = ASSETS / "render_cache"
OUT = ROOT / "outputs" / "玄女科技BP_前9页视觉返工_v0.3.pptx"


class C:
    bg = RGBColor(255, 255, 255)
    off = RGBColor(248, 250, 252)
    ink = RGBColor(15, 23, 42)
    body = RGBColor(51, 65, 85)
    muted = RGBColor(100, 116, 139)
    line = RGBColor(226, 232, 240)
    blue = RGBColor(37, 99, 235)
    sky = RGBColor(14, 165, 233)
    green = RGBColor(16, 185, 129)
    mint = RGBColor(236, 253, 245)
    pale_blue = RGBColor(239, 246, 255)
    purple = RGBColor(124, 58, 237)
    pale_purple = RGBColor(245, 243, 255)
    amber = RGBColor(245, 158, 11)
    pale_amber = RGBColor(255, 251, 235)
    dark = RGBColor(2, 6, 23)
    white = RGBColor(255, 255, 255)


FONT = "Noto Sans SC"

IMG = {
    "cover": ASSETS / "geo_embedding_cover_white_v2.png",
    "s2": USER / "patch_000339_s2_2025-04.png",
    "s2hr": USER / "patch_000339_s2_hr_2025-04.png",
    "s1": USER / "patch_000339_s1_2025-04.png",
    "landsat": USER / "patch_000339_landsat_2025-04.png",
    "dem": USER / "patch_000339_dem.png",
    "worldcover": USER / "patch_000339_worldcover.png",
    "dynamic": USER / "patch_000339_dynamic_world_2025-04_detail.png",
    "construction": USER / "patch_000217_construction_2025-08_vs_2025-09_detail.png",
    "land": USER / "patch_000217_land_conversion_2025-08_vs_2025-09_detail.png",
    "building": USER / "patch_000354_building_extraction_2025-04_detail.png",
    "cloud": USER / "patch_000146_2025-06_vs_2025-08_detail.png",
    "semantic": USER / "harbin_semantic_preset_frame.png",
    "migration": USER / "harbin_migration_frame.png",
    "migration_gif": USER / "harbin_migration_animated.gif",
    "persona_gov": ASSETS / "persona_government_v2.png",
    "persona_biz": ASSETS / "persona_enterprise_v2.png",
    "persona_uni": ASSETS / "persona_university_v2.png",
    "same_land_models": ASSETS / "same_land_multi_satellite_models.png",
    "downstream_harbin_embedding": PPT_MEDIA / "image4.png",
    "downstream_harbin_task": PPT_MEDIA / "image5.png",
    "downstream_harbin_cover": PPT_MEDIA / "image7.png",
    "downstream_harbin_landuse": PPT_MEDIA / "image8.png",
    "aef_mosaic": AEF / "deepmind_aef_hero.jpg",
    "aef_vector": AEF / "deepmind_aef_4.jpg",
    "aef_carto": AEF / "carto_aef_cover.webp",
    "aef_global_embedding": E4_AEF / "slide04_img10_24c66b14.png",
    "aef_change": E4_AEF / "slide09_img04_f5b957e4.png",
    "aef_target": E4_AEF / "slide11_img01_ea269f1d.png",
    "aef_crop": E4_AEF / "slide18_img14_6da1eb9d.png",
    "og_text_search": OPENGEOSCOPE / "Text_Search.jpg",
    "space_data_boom_arrow": ASSETS / "generated" / "commercial_space_data_boom_arrow.png",
    "annotation_video": ASSETS / "video" / "custom_annotation_demo.mov",
    "annotation_poster": ASSETS / "video" / "custom_annotation_demo_poster.png",
    "yajiang_embedding": YAJIANG / "embedding_quarter.jpeg",
    "yajiang_cluster": YAJIANG / "spatial_cluster.png",
    "yajiang_elevation": YAJIANG / "elevation_regression.png",
    "yajiang_classification": YAJIANG / "land_classification.png",
    "yajiang_retrieval": YAJIANG / "embedding_retrieval.png",
    "yajiang_change": YAJIANG / "change_detection.png",
    "yajiang_risk": YAJIANG / "slope_risk.png",
}


def bg(slide, color=C.bg) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def text(
    slide,
    value: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 14,
    color=C.ink,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, fill=C.white, line=C.line, rounded=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    return shape


def line(slide, x1, y1, x2, y2, color=C.line, width=1.2):
    shape = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


def picture(slide, path: Path, x, y, w, h):
    if path.exists():
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    rect(slide, x, y, w, h, C.off)
    text(slide, "素材待补", x, y, w, h, 11, C.muted, True, PP_ALIGN.CENTER)
    return None


def crop_to_ratio(path: Path, ratio: float) -> Path:
    """Return a center-cropped copy that matches the target aspect ratio."""
    CACHE.mkdir(parents=True, exist_ok=True)
    out = CACHE / f"{path.stem}_{int(ratio * 1000)}.png"
    if out.exists():
        return out
    if not path.exists():
        return path
    im = Image.open(path).convert("RGB")
    iw, ih = im.size
    current = iw / ih
    if current > ratio:
        nw = int(ih * ratio)
        left = (iw - nw) // 2
        box = (left, 0, left + nw, ih)
    else:
        nh = int(iw / ratio)
        top = (ih - nh) // 2
        box = (0, top, iw, top + nh)
    im.crop(box).save(out)
    return out


def picture_crop(slide, path: Path, x, y, w, h):
    return picture(slide, crop_to_ratio(path, w / h), x, y, w, h)


def picture_fit(slide, path: Path, x, y, w, h):
    if not path.exists():
        return picture(slide, path, x, y, w, h)
    with Image.open(path) as im:
        iw, ih = im.size
    target = w / h
    current = iw / ih
    if current > target:
        new_w = w
        new_h = w / current
        px = x
        py = y + (h - new_h) / 2
    else:
        new_h = h
        new_w = h * current
        px = x + (w - new_w) / 2
        py = y
    return picture(slide, path, px, py, new_w, new_h)


def diamond_image(path: Path, size=360) -> Path:
    CACHE.mkdir(parents=True, exist_ok=True)
    out = CACHE / f"diamond_{path.stem}_{size}.png"
    if out.exists():
        return out
    im = Image.open(path).convert("RGB")
    iw, ih = im.size
    side = min(iw, ih)
    left = (iw - side) // 2
    top = (ih - side) // 2
    im = im.crop((left, top, left + side, top + side)).resize((size, size))
    canvas = Image.new("RGBA", (size, size), (255, 255, 255, 0))
    mask = Image.new("L", (size, size), 0)
    points = [(size // 2, 0), (size - 1, size // 2), (size // 2, size - 1), (0, size // 2)]
    from PIL import ImageDraw

    draw = ImageDraw.Draw(mask)
    draw.polygon(points, fill=255)
    canvas.paste(im, (0, 0), mask)
    border = Image.new("RGBA", (size, size), (255, 255, 255, 0))
    bdraw = ImageDraw.Draw(border)
    bdraw.line(points + [points[0]], fill=(226, 232, 240, 255), width=5)
    canvas.alpha_composite(border)
    canvas.save(out)
    return out


def title(slide, no: str, heading: str, sub: str | None = None) -> None:
    text(slide, no, 0.62, 0.42, 0.42, 0.22, 8, C.sky, True)
    text(slide, heading, 1.04, 0.28, 11.4, 0.62, 24, C.ink, True)
    if sub:
        text(slide, sub, 1.06, 0.93, 10.6, 0.28, 10, C.muted)


def chip(slide, label: str, x, y, w, color=C.blue, fill=C.pale_blue) -> None:
    rect(slide, x, y, w, 0.34, fill, color)
    text(slide, label, x + 0.05, y + 0.02, w - 0.1, 0.30, 8, color, True, PP_ALIGN.CENTER)


def claim(slide, value: str, y=6.05, color=C.ink) -> None:
    text(slide, value, 0.86, y, 11.6, 0.42, 18, color, True, PP_ALIGN.CENTER)


def arrow(slide, x1, y1, x2, y2, color=C.line, width=1.4) -> None:
    shape = line(slide, x1, y1, x2, y2, color, width)
    shape.line.end_arrowhead = True


def image_card(slide, label: str, path: Path, x, y, w, h, color=C.line) -> None:
    rect(slide, x, y, w, h, C.white, color)
    picture_crop(slide, path, x + 0.06, y + 0.06, w - 0.12, h - 0.44)
    text(slide, label, x + 0.08, y + h - 0.30, w - 0.16, 0.20, 8, C.body, True, PP_ALIGN.CENTER)


def icon_database(slide, x, y, color=C.green) -> None:
    for i in range(3):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CAN, Inches(x), Inches(y + i * 0.13), Inches(0.36), Inches(0.34))
        shape.fill.solid()
        shape.fill.fore_color.rgb = C.white
        shape.line.color.rgb = color
        shape.line.width = Pt(1.0)


def icon_layers(slide, x, y, color=C.purple) -> None:
    for i in range(3):
        rect(slide, x + i * 0.08, y + i * 0.10, 0.46, 0.26, C.white, color, rounded=False)


def icon_loop(slide, x, y, color=C.blue) -> None:
    arc = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, Inches(x), Inches(y), Inches(0.58), Inches(0.58))
    arc.line.color.rgb = color
    arc.line.width = Pt(2.0)
    text(slide, "↻", x + 0.10, y + 0.07, 0.38, 0.30, 17, color, True, PP_ALIGN.CENTER)


def pain_row(slide, icon_fn, head: str, body: str, x, y, color, fill) -> None:
    rect(slide, x, y, 5.05, 0.92, fill, color)
    icon_fn(slide, x + 0.22, y + 0.20, color)
    text(slide, head, x + 0.86, y + 0.16, 1.45, 0.20, 11, C.ink, True)
    text(slide, body, x + 2.05, y + 0.14, 2.78, 0.46, 9, C.body)


def soft_icon(slide, kind: str, x, y, color, fill) -> None:
    bg_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.74), Inches(0.74))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = fill
    bg_shape.line.fill.background()
    cx = x + 0.37
    if kind == "repeat":
        for i in range(3):
            rect(slide, x + 0.20 + i * 0.06, y + 0.22 + i * 0.06, 0.30, 0.18, C.white, color, rounded=False)
        text(slide, "↻", x + 0.18, y + 0.34, 0.38, 0.20, 11, color, True, PP_ALIGN.CENTER)
    elif kind == "compute":
        rect(slide, x + 0.23, y + 0.23, 0.28, 0.28, C.white, color, rounded=False)
        for dx, dy in [(0.14, 0.16), (0.56, 0.16), (0.14, 0.56), (0.56, 0.56)]:
            dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + dx), Inches(y + dy), Inches(0.05), Inches(0.05))
            dot.fill.solid()
            dot.fill.fore_color.rgb = C.amber
            dot.line.fill.background()
    else:
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx - 0.07), Inches(y + 0.30), Inches(0.14), Inches(0.14))
        node.fill.solid()
        node.fill.fore_color.rgb = color
        node.line.fill.background()
        for px, py in [(x + 0.20, y + 0.18), (x + 0.54, y + 0.18), (x + 0.20, y + 0.56), (x + 0.54, y + 0.56)]:
            line(slide, cx, y + 0.37, px, py, color, 0.8)
            small = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(px - 0.035), Inches(py - 0.035), Inches(0.07), Inches(0.07))
            small.fill.solid()
            small.fill.fore_color.rgb = C.white
            small.line.color.rgb = color


def insight_row(slide, kind: str, head: str, body: str, x, y, color, fill) -> None:
    soft_icon(slide, kind, x, y + 0.05, color, fill)
    text(slide, head, x + 0.95, y + 0.05, 4.15, 0.24, 14, C.ink, True)
    text(slide, body, x + 0.95, y + 0.45, 4.05, 0.38, 9, C.body)
    line(slide, x + 0.95, y + 0.98, x + 4.92, y + 0.98, C.line, 0.8)


def flat_avatar(slide, x, y, role: str, color, accent) -> None:
    rect(slide, x, y, 3.55, 3.0, C.white, C.line)
    # soft background
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.74), Inches(y + 0.28), Inches(2.05), Inches(2.05))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.fill.transparency = 18
    circle.line.fill.background()
    # head and body
    head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 1.43), Inches(y + 0.62), Inches(0.66), Inches(0.66))
    head.fill.solid()
    head.fill.fore_color.rgb = RGBColor(248, 214, 190)
    head.line.fill.background()
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 1.12), Inches(y + 1.36), Inches(1.28), Inches(0.94))
    body.fill.solid()
    body.fill.fore_color.rgb = accent
    body.line.fill.background()
    # role props: badge / laptop / book
    if role == "gov":
        rect(slide, x + 1.50, y + 1.55, 0.24, 0.30, C.white, C.white)
        text(slide, "证", x + 1.50, y + 1.56, 0.24, 0.16, 6, accent, True, PP_ALIGN.CENTER)
        rect(slide, x + 2.28, y + 1.75, 0.50, 0.62, C.white, C.line)
        line(slide, x + 2.36, y + 1.94, x + 2.70, y + 1.94, C.line, 1)
    elif role == "biz":
        rect(slide, x + 2.28, y + 1.72, 0.62, 0.42, C.white, C.line)
        line(slide, x + 2.25, y + 2.18, x + 2.93, y + 2.18, C.line, 1.2)
        rect(slide, x + 0.64, y + 1.88, 0.48, 0.30, C.white, C.line)
    else:
        rect(slide, x + 2.25, y + 1.72, 0.56, 0.44, C.white, C.line)
        line(slide, x + 2.53, y + 1.72, x + 2.53, y + 2.16, C.line, 1)
        text(slide, "()", x + 1.37, y + 0.82, 0.80, 0.18, 11, C.ink, True, PP_ALIGN.CENTER)


def persona_card(slide, path: Path, x, y, w, h, fill, color) -> None:
    rect(slide, x, y, w, h, fill, color)
    picture_crop(slide, path, x + 0.42, y + 0.22, w - 0.84, h - 0.28)


def diamond_stack(slide, x, y, items: list[tuple[str, Path | None]]) -> None:
    for i, (label, path) in enumerate(items):
        dy = y + i * 0.47
        if path is None:
            text(slide, "…", x + 0.20, dy + 0.24, 0.48, 0.22, 18, C.muted, True, PP_ALIGN.CENTER)
        else:
            picture(slide, diamond_image(path), x, dy, 0.88, 0.88)
        if label:
            text(slide, label, x + 0.92, dy + 0.28, 1.22, 0.18, 8, C.body, True)


def task_tile(slide, label: str, path: Path | None, x, y, w=1.28, h=0.96) -> None:
    if path is None:
        text(slide, "…", x, y + 0.27, w, 0.26, 22, C.muted, True, PP_ALIGN.CENTER)
    else:
        picture_fit(slide, path, x, y, w, h)
    if label:
        text(slide, label, x, y + h + 0.08, w, 0.16, 7, C.body, True, PP_ALIGN.CENTER)


def evidence_panel(slide, path: Path, heading: str, source: str, x, y, w, h) -> None:
    rect(slide, x, y, w, h, C.white, C.line)
    picture_crop(slide, path, x + 0.06, y + 0.06, w - 0.12, h - 0.44)
    text(slide, heading, x + 0.12, y + h - 0.32, w - 0.24, 0.16, 8, C.ink, True)
    text(slide, source, x + 0.12, y + h - 0.13, w - 0.24, 0.10, 5, C.muted)


def positioning_row(slide, layer: str, players: str, role: str, x, y, w, fill, line_color, highlight=False) -> None:
    rect(slide, x, y, w, 0.70, fill, line_color)
    text(slide, layer, x + 0.18, y + 0.13, 1.45, 0.18, 9, line_color if highlight else C.ink, True)
    text(slide, players, x + 1.70, y + 0.10, 2.24, 0.20, 8, C.body, True)
    text(slide, role, x + 4.02, y + 0.11, w - 4.20, 0.24, 8, C.body)


def metric(slide, number: str, label: str, x, y, w, color, fill) -> None:
    rect(slide, x, y, w, 1.0, fill, color)
    text(slide, number, x + 0.18, y + 0.13, w - 0.36, 0.30, 18, color, True, PP_ALIGN.CENTER)
    text(slide, label, x + 0.16, y + 0.58, w - 0.32, 0.22, 8, C.body, align=PP_ALIGN.CENTER)


def small_metric(slide, number: str, label: str, x, y, color, fill) -> None:
    rect(slide, x, y, 1.18, 0.62, fill, color)
    text(slide, number, x + 0.08, y + 0.10, 1.02, 0.18, 11, color, True, PP_ALIGN.CENTER)
    text(slide, label, x + 0.08, y + 0.36, 1.02, 0.14, 6, C.body, align=PP_ALIGN.CENTER)


def aef_task(slide, label: str, path: Path, x, y) -> None:
    rect(slide, x, y, 3.58, 1.20, C.white, C.line)
    picture_crop(slide, path, x + 0.06, y + 0.06, 3.46, 0.84)
    text(slide, label, x + 0.10, y + 0.97, 3.38, 0.12, 7, C.body, True, PP_ALIGN.CENTER)


def compare_row(slide, item: str, aef: str, xuannv: str, y: float, highlight=False) -> None:
    fill = C.pale_blue if highlight else C.white
    rect(slide, 6.28, y, 1.20, 0.52, fill, C.line)
    rect(slide, 7.48, y, 2.22, 0.52, fill, C.line)
    rect(slide, 9.70, y, 2.70, 0.52, C.mint if highlight else C.white, C.green if highlight else C.line)
    text(slide, item, 6.42, y + 0.15, 0.92, 0.14, 7, C.ink, True, PP_ALIGN.CENTER)
    text(slide, aef, 7.60, y + 0.10, 1.98, 0.24, 7, C.body, True, PP_ALIGN.CENTER)
    text(slide, xuannv, 9.84, y + 0.09, 2.42, 0.26, 7, C.ink if highlight else C.body, True, PP_ALIGN.CENTER)


def difference_card(slide, head: str, aef: str, xuannv: str, x: float, y: float, color, fill) -> None:
    rect(slide, x, y, 2.82, 1.34, C.white, C.line)
    text(slide, head, x + 0.18, y + 0.15, 2.46, 0.18, 11, C.ink, True, PP_ALIGN.CENTER)
    line(slide, x + 0.28, y + 0.48, x + 2.54, y + 0.48, C.line, 0.7)
    text(slide, "AlphaEarth", x + 0.22, y + 0.63, 1.02, 0.14, 6, C.muted, True, PP_ALIGN.CENTER)
    text(slide, "玄女", x + 1.58, y + 0.63, 0.92, 0.14, 6, color, True, PP_ALIGN.CENTER)
    text(slide, aef, x + 0.16, y + 0.86, 1.18, 0.20, 8, C.body, True, PP_ALIGN.CENTER)
    text(slide, xuannv, x + 1.46, y + 0.82, 1.18, 0.26, 9, color, True, PP_ALIGN.CENTER)
    rect(slide, x + 1.38, y + 0.72, 0.02, 0.42, fill, fill, rounded=False)


def step_chain(slide, label: str, steps: list[str], x: float, y: float, color, fill) -> None:
    text(slide, label, x, y, 1.10, 0.16, 9, color, True)
    start = x
    gap = 0.12
    box_w = (5.14 - gap * (len(steps) - 1)) / len(steps)
    for i, step in enumerate(steps):
        sx = start + i * (box_w + gap)
        rect(slide, sx, y + 0.28, box_w, 0.42, fill, color)
        text(slide, step, sx + 0.03, y + 0.40, box_w - 0.06, 0.12, 6, color, True, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arrow(slide, sx + box_w + 0.02, y + 0.49, sx + box_w + gap - 0.02, y + 0.49, C.line, 0.7)


def pain_card(slide, idx: str, head: str, body: str, x: float, y: float, color, fill) -> None:
    rect(slide, x, y, 5.18, 1.10, C.white, C.line)
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.20), Inches(y + 0.23), Inches(0.56), Inches(0.56))
    badge.fill.solid()
    badge.fill.fore_color.rgb = fill
    badge.line.color.rgb = color
    badge.line.width = Pt(1.0)
    text(slide, idx, x + 0.20, y + 0.38, 0.56, 0.14, 11, color, True, PP_ALIGN.CENTER)
    text(slide, head, x + 0.96, y + 0.19, 1.46, 0.20, 14, C.ink, True)
    text(slide, body, x + 2.36, y + 0.17, 2.58, 0.44, 9, C.body)


def yajiang_task(slide, label: str, path: Path, x: float, y: float, w: float = 1.52, h: float = 1.03) -> None:
    picture_fit(slide, path, x, y, w, h)
    text(slide, label, x, y + h + 0.06, w, 0.12, 6, C.body, True, PP_ALIGN.CENTER)


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Cover
    s = prs.slides.add_slide(blank)
    bg(s)
    picture(s, IMG["cover"], 0, 0, 13.333, 7.5)
    text(s, "玄女科技", 0.82, 1.42, 4.5, 0.42, 24, C.ink, True)
    text(s, "用地理嵌入\n赋能遥感智能应用", 0.82, 2.08, 5.45, 1.15, 31, C.ink, True)
    text(s, "地球观测 - 一次表征 - 多任务复用", 0.86, 3.64, 5.4, 0.32, 15, C.body)
    text(s, "做中国的遥感数据通用嵌入底座", 0.86, 4.10, 5.4, 0.30, 13, C.blue, True)

    # 2. Company intro
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "02", "企业简介")
    text(
        s,
        "北京玄女科技有限公司是一家专注于地球空间智能的人工智能基础底座提供商，致力于打造遥感数据通用嵌入底座。",
        1.04,
        1.13,
        11.05,
        0.34,
        15,
        C.ink,
        True,
    )
    text(
        s,
        "公司以遥感领域的大规模预训练模型为核心，为城市治理、自然资源、应急灾害、农业生态等场景提供全球范围内可复用的像素级表征能力，推动遥感智能从“定制化开发”走向“底座赋能”。",
        1.04,
        1.62,
        11.10,
        0.58,
        11,
        C.body,
    )
    text(s, "多元观测", 1.34, 2.66, 1.40, 0.22, 11, C.ink, True, PP_ALIGN.CENTER)
    text(s, "地理嵌入", 5.72, 2.66, 1.40, 0.22, 11, C.ink, True, PP_ALIGN.CENTER)
    text(s, "下游任务", 9.86, 2.66, 1.40, 0.22, 11, C.ink, True, PP_ALIGN.CENTER)
    diamond_stack(
        s,
        1.10,
        3.02,
        [
            ("哨兵二号光学", IMG["s2"]),
            ("雷达影像", IMG["s1"]),
            ("高分光学", IMG["s2hr"]),
            ("", None),
            ("高程数据", IMG["dem"]),
        ],
    )
    text(s, "→", 3.52, 4.03, 0.58, 0.34, 24, C.muted, True, PP_ALIGN.CENTER)
    picture_fit(s, IMG["migration_gif"], 4.46, 3.05, 3.20, 2.72)
    text(s, "→", 8.02, 4.03, 0.58, 0.34, 24, C.muted, True, PP_ALIGN.CENTER)
    task_tile(s, "土地分类", IMG["downstream_harbin_task"], 8.88, 3.08, 1.28, 0.98)
    task_tile(s, "路网分割", IMG["downstream_harbin_cover"], 10.72, 3.08, 1.28, 0.98)
    task_tile(s, "土地利用", IMG["downstream_harbin_landuse"], 8.88, 4.54, 1.28, 0.98)
    task_tile(s, "", None, 10.72, 4.54, 1.28, 0.98)
    text(s, "玄女底座融合多元观测数据，生成统一地理嵌入，支撑不同下游任务复用。", 1.58, 5.92, 10.20, 0.22, 10, C.muted, align=PP_ALIGN.CENTER)
    claim(s, "一次融合生成统一嵌入，多任务持续复用。", 6.54, C.blue)

    # 3. Same place, multiple observations
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "03", "我们为什么要做这样一件事情")
    picture_fit(s, IMG["same_land_models"], 0.72, 1.22, 6.35, 4.62)
    text(s, "不同地理模型都在对同一地物的多源图像重复编码、解码，造成大量算力与工程成本浪费。", 0.98, 5.94, 5.82, 0.22, 8, C.muted, align=PP_ALIGN.CENTER)
    line(s, 7.18, 1.42, 7.18, 5.76, C.line, 0.9)
    insight_row(s, "repeat", "同一批观测，被反复训练", "气象、分割、变化检测模型，各自重建底层特征。", 7.55, 1.55, C.blue, C.pale_blue)
    insight_row(s, "compute", "成本耗在重复编码", "数据处理、标注、训练无法复用，任务越多成本越高。", 7.55, 3.02, C.green, C.mint)
    insight_row(s, "embedding", "缺少通用嵌入层", "行业需要统一地理表征，连接多源观测与下游应用。", 7.55, 4.49, C.purple, C.pale_purple)
    claim(s, "从每个任务重训，走向一次嵌入、多任务复用。", 6.22, C.blue)

    # 4. Benchmark and timing
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "04", "地理嵌入方向，我们走在前列")
    text(s, "AlphaEarth Foundations 已经把遥感从“影像处理”推向“统一地球表征”。", 1.04, 1.02, 10.8, 0.24, 10, C.muted)

    picture_crop(s, IMG["aef_global_embedding"], 0.86, 1.38, 6.08, 3.28)
    text(s, "AlphaEarth 全球嵌入可视化：把多源地球观测压缩为可计算的语义空间。", 0.98, 4.82, 5.84, 0.18, 8, C.body, True, PP_ALIGN.CENTER)

    line(s, 7.18, 1.40, 7.18, 4.82, C.line, 0.8)
    text(s, "最前列的 AlphaEarth 证明了什么", 7.48, 1.42, 4.45, 0.26, 15, C.ink, True)
    text(
        s,
        "不是再训练一个单任务模型，而是把光学、雷达、Landsat、地形与生态观测融合成年度像素级嵌入。",
        7.48,
        1.86,
        4.35,
        0.48,
        9,
        C.body,
    )
    text(s, "1. 数据形态：从多源影像变成 64 维向量", 7.48, 2.60, 4.26, 0.20, 9, C.ink, True)
    text(s, "2. 使用方式：从专家流程变成直接检索、分类、回归", 7.48, 3.02, 4.26, 0.20, 9, C.ink, True)
    text(s, "3. 行业意义：地理嵌入正在成为遥感应用的新入口", 7.48, 3.44, 4.26, 0.20, 9, C.ink, True)
    small_metric(s, "10 m", "像素级表征", 7.48, 4.10, C.blue, C.pale_blue)
    small_metric(s, "64 维", "嵌入向量", 8.70, 4.10, C.green, C.mint)
    small_metric(s, "1.4 万亿+", "年度足迹", 9.92, 4.10, C.purple, C.pale_purple)
    small_metric(s, "2017-2024", "年度数据", 11.14, 4.10, C.amber, C.pale_amber)
    text(s, "数据来源：Google Earth Engine Satellite Embedding 数据集", 7.52, 4.82, 4.36, 0.12, 5, C.muted)

    text(s, "一个嵌入，已经覆盖多类下游任务", 0.90, 5.14, 5.00, 0.18, 10, C.ink, True)
    aef_task(s, "相似性识别 / 变化检测", IMG["aef_change"], 0.86, 5.48)
    aef_task(s, "目标检索 / 光伏识别", IMG["aef_target"], 4.88, 5.48)
    aef_task(s, "土地分类 / 农用地识别", IMG["aef_crop"], 8.90, 5.48)
    claim(s, "AlphaEarth 验证了方向：地理嵌入会成为遥感应用的基础层。", 6.82, C.blue)

    # 5. Xuannv versus AlphaEarth
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "05", "对标 AlphaEarth：玄女的差异化路径")
    text(s, "国内较早系统性推进地理嵌入底座的团队，面向中国遥感数据与政企业务场景做高频、高分辨率、可部署的基础设施。", 1.04, 1.02, 10.70, 0.24, 10, C.muted)

    text(s, "团队具备全球数据嵌入项目开发经验", 0.88, 1.42, 4.85, 0.26, 15, C.ink, True)
    picture_crop(s, IMG["og_text_search"], 0.86, 1.86, 4.90, 2.54)
    text(s, "OpenGeoScope / EarthEmbeddingExplorer", 0.98, 4.58, 4.66, 0.18, 9, C.ink, True, PP_ALIGN.CENTER)
    text(
        s,
        "团队此前参与开发的全球卫星影像嵌入检索项目，与欧空局 ESA PhiLab / MajorTOM 生态合作，获得国际社区广泛关注与好评。",
        0.98,
        4.92,
        4.66,
        0.58,
        9,
        C.body,
        align=PP_ALIGN.CENTER,
    )
    rect(s, 1.10, 5.86, 4.44, 0.42, C.pale_blue, C.blue)
    text(s, "国际项目经验，是玄女底座能力的早期验证。", 1.36, 5.96, 3.92, 0.14, 9, C.blue, True, PP_ALIGN.CENTER)

    line(s, 6.02, 1.40, 6.02, 6.45, C.line, 0.8)
    text(s, "差异化定位：面向中国场景的高频高分辨率地理嵌入", 6.28, 1.42, 5.95, 0.30, 14, C.ink, True)
    difference_card(s, "数据体系", "全球公开数据", "加入中国遥感数据", 6.28, 2.04, C.green, C.mint)
    difference_card(s, "空间粒度", "10 米", "最高 0.5 米", 9.48, 2.04, C.blue, C.pale_blue)
    difference_card(s, "时间频率", "一年一景", "最快 11 天", 6.28, 3.76, C.purple, C.pale_purple)
    difference_card(s, "业务部署", "模型不开源", "可面向政企部署", 9.48, 3.76, C.green, C.mint)
    rect(s, 6.40, 5.72, 5.76, 0.58, RGBColor(255, 251, 235), C.amber)
    text(s, "阶段性结论：变化检测实验优于 AlphaEarth；目标识别仍需继续提升训练数据规模。", 6.70, 5.90, 5.16, 0.18, 10, C.amber, True, PP_ALIGN.CENTER)
    text(s, "GitHub: github.com/OpenGeoScope/EarthEmbeddingExplorer", 8.02, 6.40, 4.12, 0.14, 6, C.blue, True, PP_ALIGN.RIGHT)
    claim(s, "AlphaEarth 证明方向，玄女把方向落到中国高分辨率、高频更新、可部署的业务场景。", 6.82, C.blue)

    # 6. Commercial space window
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "06", "遥感应用正在迎来它的 GPT 时刻")
    text(s, "商业航天快速发展打开数据供给，应用层需要新的基础设施。", 1.04, 1.02, 10.70, 0.24, 10, C.muted)

    picture_fit(s, IMG["space_data_boom_arrow"], 0.58, 1.28, 5.78, 3.12)
    text(s, "随着商业航天的发展，遥感数据迎来自己的井喷时刻。", 3.12, 4.12, 2.72, 0.24, 10, C.blue, True, PP_ALIGN.CENTER)

    rect(s, 0.86, 4.70, 5.42, 1.55, C.white, C.line)
    text(s, "地理嵌入：遥感数据的通用表征单元", 1.10, 4.92, 4.92, 0.18, 13, C.ink, True, PP_ALIGN.CENTER)
    text(s, "我们为每个 10m x 10m 像素生成地理表征，把同一地物的多源、多时相观测压缩为统一嵌入。", 1.08, 5.30, 4.95, 0.22, 9, C.body, align=PP_ALIGN.CENTER)
    text(s, "这类似 GPT 将文本转成 token：先形成可复用的基础表征，再支撑分类、变化检测、预测和问答。", 1.08, 5.76, 4.95, 0.22, 10, C.blue, True, PP_ALIGN.CENTER)

    line(s, 6.58, 1.32, 6.58, 6.30, C.line, 0.8)
    text(s, "应用层：从项目制训练，到嵌入驱动的少样本适配", 6.92, 1.34, 5.42, 0.26, 14, C.ink, True)
    step_chain(s, "传统实现", ["数据下载", "数据预处理", "人工标注", "模型训练", "交付"], 6.92, 1.82, C.amber, RGBColor(255, 251, 235))
    text(s, "每个下游任务都要从头来一遍，工程链条长、标注成本高。", 6.98, 2.62, 5.02, 0.16, 8, C.body, True, PP_ALIGN.CENTER)
    step_chain(s, "玄女实现", ["统一嵌入", "少量标注", "任务头", "多任务复用"], 6.92, 2.90, C.green, C.mint)
    text(s, "底层表征可复用，只需面向新任务适配任务头。", 6.98, 3.70, 5.02, 0.16, 8, C.body, True, PP_ALIGN.CENTER)
    movie = s.shapes.add_movie(
        str(IMG["annotation_video"]),
        Inches(7.04),
        Inches(4.02),
        Inches(5.06),
        Inches(2.18),
        poster_frame_image=str(IMG["annotation_poster"]),
        mime_type="video/quicktime",
    )
    movie.line.color.rgb = C.line
    movie.line.width = Pt(0.8)
    text(s, "自定义标注与任务头训练展示（点击播放）", 7.16, 6.34, 4.82, 0.18, 9, C.body, True, PP_ALIGN.CENTER)
    claim(s, "当遥感数据像文本一样持续增长，地理嵌入就是遥感应用规模化的入口。", 6.82, C.blue)

    # 7. Three industry frictions
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "07", "全域底座到底解决什么痛点")
    text(s, "遥感项目真正难的不是做一个模型，而是让同一片区域的数据、表征和任务结果持续复用。", 1.04, 1.02, 10.70, 0.24, 10, C.muted)

    text(s, "全域底座解决的不是单点模型问题，而是规模化复用问题", 0.88, 1.42, 5.15, 0.24, 14, C.ink, True)
    pain_card(s, "01", "数据难统一", "多源、多时相、多分辨率数据分散在不同流程里，难以形成统一表征。", 0.86, 1.90, C.blue, C.pale_blue)
    pain_card(s, "02", "任务难复用", "每做一个下游任务，都要重新处理数据、重新标注、重新训练模型。", 0.86, 3.18, C.amber, C.pale_amber)
    pain_card(s, "03", "能力难沉淀", "项目交付后通常只留下结果图，无法把区域知识沉淀成持续调用的能力。", 0.86, 4.46, C.green, C.mint)
    rect(s, 1.08, 5.76, 4.66, 0.50, C.pale_blue, C.blue)
    text(s, "全域底座的价值：把区域数据先沉淀成统一嵌入，再支撑多任务调用。", 1.28, 5.90, 4.26, 0.14, 9, C.blue, True, PP_ALIGN.CENTER)

    line(s, 6.22, 1.42, 6.22, 6.22, C.line, 0.8)
    text(s, "雅江案例：同一套嵌入数据，支撑多类下游任务", 6.58, 1.42, 5.55, 0.26, 15, C.ink, True)
    picture_fit(s, IMG["yajiang_retrieval"], 6.60, 1.92, 2.38, 1.26)
    picture_fit(s, IMG["yajiang_embedding"], 9.30, 1.78, 2.18, 1.54)
    text(s, "雅江嵌入区域范围", 6.62, 3.30, 2.34, 0.12, 7, C.body, True, PP_ALIGN.CENTER)
    text(s, "23-26 年季度嵌入数据可视化", 9.30, 3.30, 2.18, 0.12, 7, C.body, True, PP_ALIGN.CENTER)
    text(s, "利用嵌入数据集做下游任务", 7.78, 3.62, 2.62, 0.18, 10, C.green, True, PP_ALIGN.CENTER)
    yajiang_task(s, "空间聚类", IMG["yajiang_cluster"], 6.42, 3.98, 1.70, 0.72)
    yajiang_task(s, "高程回归", IMG["yajiang_elevation"], 8.38, 3.90, 1.20, 0.88)
    yajiang_task(s, "地物分类", IMG["yajiang_classification"], 10.06, 3.96, 1.76, 0.78)
    yajiang_task(s, "嵌入数据集检索", IMG["yajiang_retrieval"], 6.52, 5.10, 1.52, 0.78)
    yajiang_task(s, "变化检测", IMG["yajiang_change"], 8.46, 5.00, 1.18, 0.90)
    yajiang_task(s, "坡度风险预测", IMG["yajiang_risk"], 10.30, 5.00, 1.18, 0.90)
    claim(s, "全域底座不是多做几个模型，而是让一个区域嵌入持续服务多个业务问题。", 6.82, C.blue)

    # 8. Angel users
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "08", "天使用户：政府、企业、高校，痛点都指向同一件事")
    users = [
        ("gov", "政府部门", "变化图斑太多", "需要证据链、优先级、核查闭环", C.pale_blue, C.blue),
        ("biz", "遥感企业", "项目越多越重", "需要复用底座、少样本适配、降低人力", C.mint, C.green),
        ("edu", "高校团队", "数据工程太重", "需要现成嵌入、预标注、难例检索", C.pale_purple, C.purple),
    ]
    for i, (role, name, pain, value, fill, color) in enumerate(users):
        x = 0.90 + i * 4.16
        path = {"gov": IMG["persona_gov"], "biz": IMG["persona_biz"], "edu": IMG["persona_uni"]}[role]
        persona_card(s, path, x, 1.40, 3.55, 3.0, fill, color)
        text(s, name, x + 0.16, 4.58, 3.23, 0.24, 16, C.ink, True, PP_ALIGN.CENTER)
        text(s, pain, x + 0.22, 5.08, 3.10, 0.20, 11, color, True, PP_ALIGN.CENTER)
        text(s, value, x + 0.34, 5.48, 2.86, 0.34, 9, C.body, align=PP_ALIGN.CENTER)
    claim(s, "三类用户表面不同，本质都是：不要每个任务重做一次遥感工程。", 6.23, C.ink)

    # 9. HMW and token analogy
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "09", "我们的答案：把地球观测转成地理智能底座")
    text(
        s,
        "我们如何把爆发式增长的地球观测数据转化为可复用的地理智能底座，\n让不同用户不再为每个遥感任务重复采购、标注、建模和计算？",
        1.04,
        1.30,
        11.25,
        0.96,
        24,
        C.ink,
        True,
        PP_ALIGN.CENTER,
    )
    rect(s, 1.30, 3.22, 4.75, 1.38, C.pale_blue, C.blue)
    text(s, "文本智能", 1.62, 3.48, 4.10, 0.24, 14, C.blue, True, PP_ALIGN.CENTER)
    text(s, "文本 → 语义单元 → 理解 / 检索 / 生成", 1.62, 3.94, 4.10, 0.24, 11, C.body, True, PP_ALIGN.CENTER)
    rect(s, 7.28, 3.22, 4.75, 1.38, C.mint, C.green)
    text(s, "地球智能", 7.60, 3.48, 4.10, 0.24, 14, C.green, True, PP_ALIGN.CENTER)
    text(s, "地球观测 → 地理嵌入 → 监测 / 核查 / 预警 / 决策", 7.60, 3.94, 4.10, 0.24, 10, C.body, True, PP_ALIGN.CENTER)
    arrow(s, 6.18, 3.92, 7.05, 3.92, C.line, 1.4)
    text(
        s,
        "大模型先把文本转成语义单元，才能进行理解、检索、生成；\n玄女先把地球观测转成地理嵌入，才能进行监测、核查、预警和决策。",
        1.16,
        5.45,
        11.0,
        0.58,
        15,
        C.blue,
        True,
        PP_ALIGN.CENTER,
    )

    return prs


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
