"""Generate the first revised BP pages requested by the user.

The user added a company-introduction slide between page 1 and page 2, so this
deck has 9 slides: cover + intro + the user's revised pages 2-8.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
OUT = ROOT / "outputs" / "玄女科技BP_前9页_含企业简介_v0.1.pptx"


class C:
    bg = RGBColor(255, 255, 255)
    off = RGBColor(247, 250, 252)
    ink = RGBColor(17, 24, 39)
    muted = RGBColor(100, 116, 139)
    line = RGBColor(226, 232, 240)
    blue = RGBColor(82, 148, 226)
    green = RGBColor(72, 187, 143)
    purple = RGBColor(151, 132, 236)
    cyan = RGBColor(92, 200, 220)
    pale_blue = RGBColor(235, 245, 255)
    pale_green = RGBColor(235, 250, 243)
    pale_purple = RGBColor(244, 241, 255)
    white = RGBColor(255, 255, 255)


FONT = "Noto Sans SC"


IMG = {
    "cover": ASSETS / "geo_embedding_cover_pastel.png",
    "grid": ASSETS / "grid_satellite_embedding.png",
    "embed": ASSETS / "harbin_embedding_pca.png",
    "optical": ASSETS / "harbin_highres_optical.png",
    "sar": ASSETS / "harbin_s1_sar.png",
    "prediction": ASSETS / "harbin_prediction.png",
    "worldcover": ASSETS / "harbin_worldcover.png",
    "old_earth": ASSETS / "old_bp_media" / "image1.png",
    "old_optical": ASSETS / "old_bp_media" / "image15.jpeg",
    "old_sar": ASSETS / "old_bp_media" / "image16.png",
    "old_embed": ASSETS / "old_bp_media" / "image17.png",
    "government": ASSETS / "persona_government.png",
    "enterprise": ASSETS / "persona_enterprise.png",
    "university": ASSETS / "persona_university.png",
}


def bg(slide, color=C.bg):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def text(slide, s, x, y, w, h, size=14, color=C.ink, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = s
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return shape


def rect(slide, x, y, w, h, fill=C.white, line=C.line, rounded=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(0.8)
    return shp


def line(slide, x1, y1, x2, y2, color=C.line, width=1.2):
    shp = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shp.line.color.rgb = color
    shp.line.width = Pt(width)
    return shp


def picture(slide, key: str, x, y, w, h):
    path = IMG[key]
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        rect(slide, x, y, w, h, C.off)
        text(slide, "素材待补", x, y + h / 2 - 0.12, w, 0.24, 11, C.muted, True, PP_ALIGN.CENTER)


def title(slide, no: str, heading: str, sub: str = ""):
    text(slide, no, 0.62, 0.38, 0.42, 0.20, 8, C.blue, True)
    text(slide, heading, 1.08, 0.30, 10.8, 0.55, 23, C.ink, True)
    if sub:
        text(slide, sub, 1.10, 0.93, 10.7, 0.28, 10, C.muted)


def footer(slide, no: int):
    text(slide, "玄女科技商业计划书｜前置叙事 v0.1", 0.62, 7.15, 4.8, 0.18, 7, C.muted)
    text(slide, f"{no:02d}", 12.1, 7.12, 0.45, 0.18, 8, C.muted, align=PP_ALIGN.RIGHT)


def pill(slide, s, x, y, w, color=C.blue):
    rect(slide, x, y, w, 0.32, C.off, color)
    text(slide, s, x + 0.08, y + 0.07, w - 0.16, 0.14, 7, color, True, PP_ALIGN.CENTER)


def card(slide, head, body, x, y, w, h, color=C.blue, fill=C.white):
    rect(slide, x, y, w, h, fill)
    text(slide, head, x + 0.2, y + 0.16, w - 0.4, 0.22, 11, C.ink, True)
    text(slide, body, x + 0.2, y + 0.52, w - 0.4, h - 0.62, 8, C.muted)
    line(slide, x + 0.2, y + h - 0.14, x + w - 0.2, y + h - 0.14, color, 2.0)


def flow(slide, items, x, y, w, color=C.blue):
    gap = 0.14
    bw = (w - gap * (len(items) - 1)) / len(items)
    for i, item in enumerate(items):
        bx = x + i * (bw + gap)
        rect(slide, bx, y, bw, 0.58, C.off, color)
        text(slide, item, bx + 0.05, y + 0.18, bw - 0.1, 0.15, 8, C.ink, True, PP_ALIGN.CENTER)
        if i < len(items) - 1:
            text(slide, "→", bx + bw + 0.03, y + 0.19, 0.08, 0.15, 8, C.muted, True)


def mini_chart(slide, labels, values, x, y, w, h, color=C.blue):
    max_v = max(values)
    step = w / (len(values) * 2 + 1)
    for i, (lab, val) in enumerate(zip(labels, values)):
        bx = x + step + i * step * 2
        bh = h * val / max_v
        rect(slide, bx, y + h - bh, step, bh, color, color, rounded=False)
        text(slide, lab, bx - 0.25, y + h + 0.12, step + 0.5, 0.18, 7, C.muted, align=PP_ALIGN.CENTER)
        text(slide, str(val), bx - 0.15, y + h - bh - 0.22, step + 0.3, 0.16, 7, C.ink, True, PP_ALIGN.CENTER)


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 cover
    s = prs.slides.add_slide(blank); bg(s)
    picture(s, "cover", 5.4, 0.0, 7.95, 7.5)
    pill(s, "商业计划书", 0.86, 0.86, 1.18, C.blue)
    text(s, "玄女科技", 0.86, 1.58, 4.2, 0.48, 25, C.ink, True)
    text(s, "用地理嵌入赋能遥感智能应用", 0.86, 2.18, 5.7, 0.68, 28, C.ink, True)
    text(s, "地球观测，一次表征，多任务复用", 0.90, 3.15, 5.2, 0.30, 15, C.muted)
    text(s, "让中国拥有自己的地理智能底座", 0.90, 3.62, 5.2, 0.26, 12, C.blue, True)

    # 2 company intro
    s = prs.slides.add_slide(blank); bg(s); title(s, "02", "企业简介：玄女底座在做什么", "一句话让投资人看懂公司定位与产品形态。")
    text(s, "玄女科技是一家面向地球空间智能的 AI 基础底座公司，将多源地球观测数据转化为可复用地理嵌入，支撑变化检测、分类、预测和问答等下游任务。", 1.08, 1.28, 10.9, 0.45, 15, C.ink, True)
    flow(s, ["光学", "SAR", "Landsat", "高分", "DEM", "气象"], 0.75, 2.35, 3.2, C.blue)
    rect(s, 4.38, 1.92, 3.35, 3.35, fill=RGBColor(250, 252, 255), line=C.blue)
    picture(s, "embed", 4.55, 2.08, 3.0, 3.0)
    text(s, "地理嵌入 PCA 可视化", 4.55, 5.25, 3.0, 0.20, 8, C.muted, align=PP_ALIGN.CENTER)
    for i, (name, key) in enumerate([("施工", "prediction"), ("建筑", "optical"), ("农用地", "worldcover"), ("灾害", "old_sar")]):
        x = 8.35 + (i % 2) * 1.65
        y = 2.10 + (i // 2) * 1.75
        picture(s, key, x, y, 1.15, 1.15)
        text(s, name, x, y + 1.22, 1.15, 0.15, 7, C.muted, True, PP_ALIGN.CENTER)
    text(s, "输入源数据", 1.45, 3.52, 1.2, 0.18, 10, C.blue, True, PP_ALIGN.CENTER)
    text(s, "统一地理嵌入", 5.0, 1.72, 2.3, 0.18, 10, C.blue, True, PP_ALIGN.CENTER)
    text(s, "下游任务", 8.75, 1.72, 2.3, 0.18, 10, C.blue, True, PP_ALIGN.CENTER)
    footer(s, 2)

    # 3 why do it
    s = prs.slides.add_slide(blank); bg(s); title(s, "03", "为什么要做：同一块地表被反复编码，数据却仍是孤岛")
    picture(s, "grid", 0.55, 1.25, 6.05, 4.35)
    card(s, "重复造轮子", "气象、水文、分割、变化检测等模型，都从同一块地表影像重复做编码器和解码器。", 7.0, 1.55, 5.2, 1.05, C.blue)
    card(s, "数据孤岛", "光学、SAR、气象、水文、DEM、高分数据各自处理，难以形成统一资产。", 7.0, 2.92, 5.2, 1.05, C.green)
    card(s, "玄女底座", "联合多源观测生成统一地理嵌入，减少冗余计算，并为下游任务提供更稳定表征。", 7.0, 4.29, 5.2, 1.05, C.purple)
    text(s, "从“每个任务重新编码”到“一个底座多任务复用”。", 1.0, 6.15, 11.2, 0.30, 18, C.ink, True, PP_ALIGN.CENTER)
    footer(s, 3)

    # 4 benchmark
    s = prs.slides.add_slide(blank); bg(s); title(s, "04", "国外已有 AlphaEarth 方向，玄女国内率先做中国自己的地理智能底座")
    picture(s, "old_earth", 0.8, 1.35, 3.25, 1.85)
    card(s, "国外趋势", "AlphaEarth / Google Satellite Embedding 已经证明：把地球观测转成嵌入，是下一代地理计算入口。", 0.8, 3.55, 3.25, 1.1, C.blue)
    mini_chart(s, ["频次", "分辨率", "本土化"], [11, 10, 9], 4.9, 1.82, 2.8, 1.6, C.green)
    card(s, "玄女优势 1：频次更高", "最快可做到 11 天一次嵌入，支撑更高频变化监测。", 8.35, 1.35, 3.7, 0.86, C.green)
    card(s, "玄女优势 2：分辨率更高", "对标国外 10 米级，并面向高分光学 / SAR 等更高分模态融合。", 8.35, 2.52, 3.7, 0.86, C.blue)
    card(s, "玄女优势 3：更本土", "适配中国卫星数据、国产算力和本土政企场景。", 8.35, 3.69, 3.7, 0.86, C.purple)
    footer(s, 4)

    # 5 urgency
    s = prs.slides.add_slide(blank); bg(s); title(s, "05", "为什么现在急需做：商业航天很热，但数据处理链条仍滞后")
    flow(s, ["卫星制造", "火箭发射", "星座运营", "数据下行", "数据处理", "智能应用"], 0.8, 1.75, 11.6, C.blue)
    rect(s, 8.42, 1.62, 3.95, 0.85, fill=RGBColor(241, 248, 255), line=C.green)
    text(s, "行业断点", 9.55, 1.42, 1.2, 0.18, 9, C.green, True, PP_ALIGN.CENTER)
    picture(s, "old_earth", 0.85, 3.0, 3.0, 1.7)
    picture(s, "old_optical", 4.2, 3.0, 3.0, 1.7)
    picture(s, "old_sar", 7.55, 3.0, 3.0, 1.7)
    card(s, "玄女切入点", "卫星越来越多，数据越来越多；真正缺的是把数据处理成可复用智能应用的底座。", 1.25, 5.35, 10.7, 0.78, C.green)
    footer(s, 5)

    # 6 why industry not developed
    s = prs.slides.add_slide(blank); bg(s); title(s, "06", "为什么遥感行业还没有真正发展起来：用不起、用不好、用不快")
    card(s, "用不起", "项目制、专家制、标注成本高，边际成本降不下来。", 0.85, 1.52, 3.45, 1.05, C.blue, C.pale_blue)
    card(s, "用不好", "跨区域泛化差，换任务、换地物、换传感器就要重做。", 4.9, 1.52, 3.45, 1.05, C.green, C.pale_green)
    card(s, "用不快", "交付周期长，从需求提出到结果交付常常错过决策窗口。", 8.95, 1.52, 3.45, 1.05, C.purple, C.pale_purple)
    flow(s, ["数据采购", "预处理", "专家建模", "特征提取", "制图交付"], 0.95, 3.35, 5.4, C.muted)
    flow(s, ["地理嵌入", "任务适配", "候选图斑", "证据链"], 7.0, 3.35, 4.8, C.green)
    text(s, "问题不在于没有卫星数据，而在于每个应用都要重新走一遍重流程。", 1.0, 5.55, 11.2, 0.38, 19, C.ink, True, PP_ALIGN.CENTER)
    footer(s, 6)

    # 7 users
    s = prs.slides.add_slide(blank); bg(s); title(s, "07", "天使用户与目标用户：政府、企业、学校，未来开放给个人")
    personas = [("政府", "government", "需要变化证据链", "图斑优先级、历史过程、核查闭环"), ("企业", "enterprise", "项目越多越重", "底座复用、少样本适配、降本增效"), ("学校", "university", "数据工程吃掉周期", "现成嵌入、预标注、难例检索")]
    for i, (name, key, pain, solve) in enumerate(personas):
        x = 0.85 + i * 4.12
        picture(s, key, x, 1.45, 3.15, 2.35)
        text(s, name, x, 4.02, 3.15, 0.22, 14, C.ink, True, PP_ALIGN.CENTER)
        text(s, pain, x + 0.18, 4.45, 2.8, 0.20, 10, C.blue, True, PP_ALIGN.CENTER)
        text(s, solve, x + 0.18, 4.85, 2.8, 0.42, 8, C.muted, align=PP_ALIGN.CENTER)
    text(s, "规模做上来后，可把地块异常提醒和证据包开放给个人用户。", 1.0, 6.1, 11.2, 0.26, 14, C.muted, align=PP_ALIGN.CENTER)
    footer(s, 7)

    # 8 HMW/token
    s = prs.slides.add_slide(blank); bg(s); title(s, "08", "How might we：把地球观测转成可复用的地理智能底座")
    text(s, "我们如何把爆发式增长的地球观测数据转化为可复用的地理智能底座，让不同用户不再为每个遥感任务重复采购、标注、建模和计算？", 1.0, 1.52, 11.2, 0.82, 25, C.ink, True, PP_ALIGN.CENTER)
    card(s, "GPT", "文本 → token → 理解 / 检索 / 生成", 1.35, 3.35, 4.5, 1.15, C.blue, C.pale_blue)
    card(s, "玄女", "地球观测 → 地理嵌入 → 变化检测 / 分类 / 预测 / 问答", 7.35, 3.35, 4.5, 1.15, C.green, C.pale_green)
    text(s, "商业航天解决“看见地球”，地理嵌入解决“理解地球”。", 1.0, 5.7, 11.2, 0.34, 18, C.blue, True, PP_ALIGN.CENTER)
    footer(s, 8)

    # 9 process/tech
    s = prs.slides.add_slide(blank); bg(s); title(s, "09", "技术与业务范式：从重复处理到统一底座调用")
    flow(s, ["遥感数据采集", "数据预处理", "建模", "特征提取", "制图", "交付"], 0.75, 1.55, 11.8, C.muted)
    text(s, "传统流程：每个项目重复走长链条", 0.9, 2.42, 4.6, 0.18, 9, C.muted, True)
    flow(s, ["多源输入", "时空对齐", "地理嵌入", "任务库", "证据链 / 报告"], 0.75, 3.35, 11.8, C.green)
    text(s, "玄女流程：一次底座，多任务复用", 0.9, 4.22, 4.6, 0.18, 9, C.green, True)
    picture(s, "old_embed", 0.95, 5.05, 2.35, 1.22)
    picture(s, "old_sar", 3.65, 5.05, 2.35, 1.22)
    picture(s, "embed", 6.35, 5.05, 2.35, 1.22)
    picture(s, "prediction", 9.05, 5.05, 2.35, 1.22)
    footer(s, 9)

    return prs


def main() -> None:
    ROOT.joinpath("outputs").mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
