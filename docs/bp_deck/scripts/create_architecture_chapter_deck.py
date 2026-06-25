"""Generate the Xuannv architecture and paradigm chapter deck.

This chapter follows the narrative plan after the case deck:
paradigm shift -> technical principle -> architecture -> product/business
architecture -> roadmap.
"""

from __future__ import annotations

from pathlib import Path

import fitz
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from create_front8_deck_v2 import C, IMG, bg, title, text, rect, line, arrow, picture_fit


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs" / "玄女科技BP_架构原理章节_v0.2.pptx"
CACHE = ROOT / "assets" / "render_cache"

HAIDIAN = Path("/root/workspace/report/assets")
YAJIANG_AGENT = Path("/root/.codex/attachments/071916f0-9ca0-4d8c-83f1-bcd404e7ea3c/e178b904e06c08479cd215ca230ff245.png")
MODEL_REPORT = Path("/root/workspace/xuannv/reports/harbin_202512_202605/pdfs")


def small_label(slide, value: str, x: float, y: float, w: float, color=C.body) -> None:
    text(slide, value, x, y, w, 0.13, 7, color, True, PP_ALIGN.CENTER)


def chip(slide, value: str, x: float, y: float, w: float, color=C.blue, fill=C.pale_blue) -> None:
    rect(slide, x, y, w, 0.32, fill, color)
    text(slide, value, x + 0.05, y + 0.085, w - 0.10, 0.08, 8, color, True, PP_ALIGN.CENTER)


def slim_node(slide, value: str, x: float, y: float, w: float, color=C.blue, fill=C.pale_blue) -> None:
    rect(slide, x, y, w, 0.48, fill, color)
    text(slide, value, x + 0.08, y + 0.145, w - 0.16, 0.10, 8, color, True, PP_ALIGN.CENTER)


def evidence_card(slide, label: str, image: Path, x: float, y: float, w: float, h: float, color=C.blue) -> None:
    picture_fit(slide, image, x, y, w, h)
    small_label(slide, label, x, y + h + 0.10, w, color)


def statement(slide, value: str, y: float, color=C.blue) -> None:
    text(slide, value, 0.86, y, 11.60, 0.28, 15, color, True, PP_ALIGN.CENTER)


def section_header(slide, value: str, x: float, y: float, color=C.blue) -> None:
    text(slide, value, x, y, 3.00, 0.20, 14, color, True)
    line(slide, x, y + 0.34, x + 1.34, y + 0.34, color, 1.0)


def render_pdf_page(pdf: Path, page: int, name: str, zoom: float = 1.2) -> Path:
    CACHE.mkdir(parents=True, exist_ok=True)
    out = CACHE / f"{name}_p{page + 1}.png"
    if out.exists():
        return out
    doc = fitz.open(pdf)
    pix = doc[page].get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
    pix.save(out)
    return out


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    haidian_patch = HAIDIAN / "compare_haidian_construction_patch_000198.png"
    haidian_chart = HAIDIAN / "aef_vs_v1.0_v2.png"
    harbin_report_cover = render_pdf_page(
        MODEL_REPORT / "00-2025年12月-2026年5月哈尔滨新区城市管理卫星遥感监测总报告.pdf",
        0,
        "arch_harbin_report",
    )

    # 19. Paradigm shift
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "19", "从案例到范式：遥感应用从项目制走向底座化")
    text(s, "前面的案例证明：同一套地理嵌入可以支撑变化检测、分类评测、智能体报告等多类真实任务。", 1.02, 1.14, 11.20, 0.28, 13, C.body, False, PP_ALIGN.CENTER)
    section_header(s, "传统遥感项目制", 0.86, 1.62, C.amber)
    for i, label in enumerate(["数据采购", "预处理", "专家标注", "单任务模型", "制图报告"]):
        x = 0.94 + i * 1.02
        slim_node(s, label, x, 2.26, 0.82, C.amber, C.pale_amber)
        if i < 4:
            arrow(s, x + 0.82, 2.50, x + 0.98, 2.50, C.amber, 1.0)
    text(s, "痛点：新任务通常意味着重新组织数据、标注、模型和交付流程。", 1.00, 3.16, 4.86, 0.22, 11, C.body, True, PP_ALIGN.CENTER)
    line(s, 6.30, 1.54, 6.30, 5.78, C.line, 0.9)
    section_header(s, "玄女底座化交付", 6.84, 1.62, C.blue)
    slim_node(s, "多源观测", 6.92, 2.18, 1.12, C.green, C.mint)
    arrow(s, 8.04, 2.42, 8.40, 2.42, C.line, 1.2)
    rect(s, 8.44, 2.02, 1.74, 0.80, C.pale_blue, C.blue)
    text(s, "DV 地理嵌入向量", 8.62, 2.30, 1.38, 0.10, 9, C.blue, True, PP_ALIGN.CENTER)
    arrow(s, 10.18, 2.42, 10.52, 2.42, C.line, 1.2)
    slim_node(s, "多任务复用", 10.56, 2.18, 1.34, C.purple, C.pale_purple)
    for label, x, color, fill in [
        ("哈尔滨报告", 6.96, C.blue, C.pale_blue),
        ("海淀评测", 8.26, C.green, C.mint),
        ("雅江智能体", 9.56, C.purple, C.pale_purple),
        ("更多任务", 10.86, C.amber, C.pale_amber),
    ]:
        chip(s, label, x, 3.16, 1.04, color, fill)
    evidence_card(s, "自动报告", harbin_report_cover, 6.98, 3.90, 1.18, 1.46, C.blue)
    evidence_card(s, "Benchmark", haidian_chart, 8.62, 3.94, 1.84, 1.18, C.green)
    evidence_card(s, "智能体界面", YAJIANG_AGENT, 10.72, 3.92, 1.34, 1.22, C.purple)
    statement(s, "玄女不是多做几个遥感项目，而是把遥感项目沉淀为可复用的地理智能基础设施。", 6.58)

    # 20. Technical paradigm
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "20", "技术范式：从“图像学 + 专家规则”到“多模态自学习”")
    text(s, "传统遥感依赖显性规则和专家经验；玄女让模型从多模态、多时序观测中学习隐性地理规律。", 1.02, 1.14, 11.20, 0.28, 13, C.body, False, PP_ALIGN.CENTER)
    section_header(s, "传统图像学", 0.86, 1.62, C.amber)
    for i, label in enumerate(["定义特征", "人工融合", "单任务训练", "换任务重做"]):
        slim_node(s, label, 0.94 + i * 1.20, 2.24, 0.94, C.amber, C.pale_amber)
        if i < 3:
            arrow(s, 1.88 + i * 1.20, 2.48, 2.08 + i * 1.20, 2.48, C.amber, 1.0)
    evidence_card(s, "影像", IMG["s2"], 1.14, 3.36, 1.08, 0.96, C.body)
    evidence_card(s, "人工标签", IMG["worldcover"], 2.52, 3.36, 1.08, 0.96, C.body)
    evidence_card(s, "单任务输出", IMG["building"], 3.90, 3.36, 1.08, 0.96, C.body)
    text(s, "特点：知识写在规则里，能力随项目结束而难以沉淀。", 1.04, 5.16, 4.74, 0.20, 10, C.body, True, PP_ALIGN.CENTER)
    line(s, 6.30, 1.54, 6.30, 5.78, C.line, 0.9)
    section_header(s, "玄女自学习", 6.84, 1.62, C.blue)
    for i, label in enumerate(["多模态对齐", "时空规律学习", "DV 地理嵌入", "任务激活"]):
        slim_node(s, label, 6.92 + i * 1.28, 2.24, 1.08, C.blue if i != 1 else C.green, C.pale_blue if i != 1 else C.mint)
        if i < 3:
            arrow(s, 8.00 + i * 1.28, 2.48, 8.20 + i * 1.28, 2.48, C.blue, 1.0)
    evidence_card(s, "高分影像", IMG["s2hr"], 7.06, 3.36, 1.08, 0.96, C.body)
    evidence_card(s, "地形", IMG["dem"], 8.44, 3.36, 1.08, 0.96, C.body)
    evidence_card(s, "语义结果", IMG["semantic"], 9.82, 3.36, 1.08, 0.96, C.body)
    text(s, "特点：知识沉淀在表征里，新增任务通过轻量适配复用底座能力。", 7.02, 5.16, 4.76, 0.20, 10, C.body, True, PP_ALIGN.CENTER)
    statement(s, "技术变化的本质：从人工解释影像，转向模型自学习地理规律。", 6.58)

    # 21. Core principle
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "21", "核心原理：从“数字地球”到“嵌入地球”")
    text(s, "数字地球解决数据存储；嵌入地球解决知识编码。玄女把同一地理位置的多源观测转成可计算的 DV 地理嵌入向量。", 1.02, 1.14, 11.20, 0.30, 13, C.body, False, PP_ALIGN.CENTER)
    section_header(s, "1 数据图层", 0.86, 1.68, C.amber)
    for i, (label, img) in enumerate([("光学", IMG["s2"]), ("SAR", IMG["s1"]), ("DEM", IMG["dem"])]):
        evidence_card(s, label, img, 0.96 + i * 1.30, 2.24, 1.00, 0.92, C.body)
    text(s, "分散、稀疏、依赖人工解释", 1.08, 3.72, 3.62, 0.16, 10, C.body, True, PP_ALIGN.CENTER)
    arrow(s, 4.98, 3.00, 5.70, 3.00, C.line, 1.4)
    section_header(s, "2 地理嵌入", 5.82, 1.68, C.blue)
    rect(s, 5.92, 2.30, 2.30, 1.06, C.pale_blue, C.blue)
    text(s, "DV 地理嵌入向量", 6.12, 2.66, 1.90, 0.12, 13, C.blue, True, PP_ALIGN.CENTER)
    for i, label in enumerate(["时间", "空间", "频谱", "语义"]):
        chip(s, label, 5.86 + i * 0.62, 3.78, 0.48, C.blue, C.white)
    text(s, "统一、致密、可比较", 5.92, 4.46, 2.30, 0.16, 10, C.body, True, PP_ALIGN.CENTER)
    arrow(s, 8.50, 3.00, 9.22, 3.00, C.line, 1.4)
    section_header(s, "3 下游调用", 9.38, 1.68, C.green)
    for i, label in enumerate(["变化检测", "地物分类", "检索问答", "报告生成"]):
        chip(s, label, 9.50 + (i % 2) * 1.12, 2.40 + (i // 2) * 0.72, 0.94, C.green, C.mint)
    text(s, "同一表征，多类任务复用", 9.48, 4.46, 2.38, 0.16, 10, C.body, True, PP_ALIGN.CENTER)
    statement(s, "GPT 先把文本转成 token；玄女先把地球观测转成地理嵌入。", 6.58)

    # 22. Model architecture
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "22", "模型架构：多模态观测进入统一地理嵌入空间")
    text(s, "玄女底座将不同传感器、不同时间、不同分辨率的数据对齐到同一套 DV 地理嵌入向量，再由任务头完成不同下游任务。", 1.02, 1.14, 11.20, 0.30, 13, C.body, False, PP_ALIGN.CENTER)
    text(s, "Input", 0.88, 1.70, 2.42, 0.18, 14, C.blue, True, PP_ALIGN.CENTER)
    for i, (label, img) in enumerate([("可见光", IMG["s2"]), ("SAR", IMG["s1"]), ("DEM", IMG["dem"]), ("多光谱", IMG["landsat"])]):
        evidence_card(s, label, img, 0.98 + (i % 2) * 1.18, 2.10 + (i // 2) * 1.28, 0.92, 0.82, C.body)
    arrow(s, 3.30, 3.10, 4.00, 3.10, C.line, 1.4)
    rect(s, 4.18, 2.12, 2.04, 1.92, C.pale_blue, C.blue)
    text(s, "时空多模态编码器", 4.36, 2.50, 1.68, 0.18, 13, C.blue, True, PP_ALIGN.CENTER)
    text(s, "时空对齐\n跨模态融合\n时序学习", 4.58, 2.94, 1.24, 0.44, 9, C.body, True, PP_ALIGN.CENTER)
    arrow(s, 6.22, 3.10, 6.92, 3.10, C.line, 1.4)
    rect(s, 7.10, 2.10, 2.00, 1.96, C.mint, C.green)
    text(s, "DV 地理嵌入向量", 7.26, 2.52, 1.68, 0.18, 13, C.green, True, PP_ALIGN.CENTER)
    text(s, "像素级 / 格网级\n可复用表征", 7.38, 2.98, 1.44, 0.36, 9, C.body, True, PP_ALIGN.CENTER)
    arrow(s, 9.10, 3.10, 9.80, 3.10, C.line, 1.4)
    text(s, "Task Heads", 9.92, 1.70, 2.46, 0.18, 14, C.blue, True, PP_ALIGN.CENTER)
    for i, label in enumerate(["LUCC", "变化检测", "高程回归", "地物分类", "检索问答", "报告生成"]):
        chip(s, label, 9.80 + (i % 2) * 1.14, 2.16 + (i // 2) * 0.66, 1.00, C.purple if i % 2 else C.blue, C.pale_purple if i % 2 else C.pale_blue)
    line(s, 0.96, 4.78, 12.20, 4.78, C.line, 0.8)
    text(s, "训练路线", 0.96, 5.20, 1.40, 0.18, 13, C.blue, True)
    for i, label in enumerate(["多源数据汇聚", "自监督预训练", "通用地理嵌入", "场景增强", "报告/API"]):
        slim_node(s, label, 2.42 + i * 1.72, 5.08, 1.28, C.blue if i != 3 else C.green, C.pale_blue if i != 3 else C.mint)
        if i < 4:
            arrow(s, 3.70 + i * 1.72, 5.32, 3.98 + i * 1.72, 5.32, C.line, 1.0)
    statement(s, "能力复用来自统一嵌入空间，而不是为每个场景重新堆模型。", 6.58)

    # 23. Product and business architecture
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "23", "产品架构：一个通用底座，支撑多类场景应用")
    text(s, "玄女把数据、模型、任务头、智能体和报告工作流封装成可调用的产品能力。", 1.02, 1.14, 11.20, 0.28, 13, C.body, False, PP_ALIGN.CENTER)
    for name, items, y, color, fill in [
        ("数据层", "遥感 / DEM / 气象 / 社会经济 / 本土标签", 1.86, C.green, C.mint),
        ("底座层", "多模态编码器 / DV 地理嵌入库 / 向量检索 / 任务头", 3.06, C.blue, C.pale_blue),
        ("应用层", "政务监测 / 城市治理 / 水电站监测 / 智能体报告 / API", 4.26, C.purple, C.pale_purple),
    ]:
        rect(s, 0.96, y, 11.28, 0.82, fill, color)
        text(s, name, 1.20, y + 0.22, 1.30, 0.16, 12, color, True)
        text(s, items, 2.56, y + 0.22, 9.28, 0.16, 10, C.ink, True, PP_ALIGN.RIGHT)
    arrow(s, 6.60, 2.70, 6.60, 3.02, C.line, 1.2)
    arrow(s, 6.60, 3.90, 6.60, 4.22, C.line, 1.2)
    text(s, "业务范式变化", 0.96, 5.48, 2.20, 0.18, 13, C.blue, True)
    for i, (old, new) in enumerate([("专家项目制", "底座调用式"), ("一次性交付", "持续服务"), ("人力扩张", "数据闭环")]):
        text(s, old, 3.10 + i * 2.88, 5.46, 1.02, 0.14, 9, C.amber, True, PP_ALIGN.CENTER)
        arrow(s, 4.10 + i * 2.88, 5.54, 4.48 + i * 2.88, 5.54, C.line, 1.0)
        text(s, new, 4.56 + i * 2.88, 5.46, 1.02, 0.14, 9, C.blue, True, PP_ALIGN.CENTER)
    statement(s, "商业价值来自底座复用：把高成本定制服务，变成可持续调用的地理智能能力。", 6.58)

    # 24. Roadmap and public benefit
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "24", "三年路径：从技术验证到地理智能普惠化")
    text(s, "未来的玄女底座不仅理解遥感影像，也将接入人口、产业、交通、用地、夜光等社会经济数据，让地理智能进入更多公共服务和商业场景。", 1.02, 1.14, 11.20, 0.30, 13, C.body, False, PP_ALIGN.CENTER)
    years = [
        ("2026", "技术验证", "中国区域 DV 地理嵌入底座验证；形成标杆案例。", "试点收入"),
        ("2027", "商业验证", "行业模型、API/License、智能体报告产品化。", "千万级目标"),
        ("2028", "规模化", "融合社会经济数据，扩展到公共服务和大众场景。", "亿元级目标"),
    ]
    for i, (yr, head, body, money) in enumerate(years):
        x = 0.92 + i * 4.02
        rect(s, x, 1.98, 3.34, 2.18, C.white, C.line)
        text(s, yr, x + 0.20, 2.20, 0.80, 0.24, 18, C.blue if i < 2 else C.green, True)
        text(s, head, x + 1.04, 2.26, 1.88, 0.16, 13, C.ink, True, PP_ALIGN.RIGHT)
        text(s, body, x + 0.24, 2.78, 2.86, 0.46, 9, C.body, True, PP_ALIGN.CENTER)
        chip(s, money, x + 0.86, 3.54, 1.56, C.blue if i < 2 else C.green, C.pale_blue if i < 2 else C.mint)
        if i < 2:
            arrow(s, x + 3.34, 3.08, x + 3.76, 3.08, C.line, 1.2)
    section_header(s, "社会经济数据接入", 0.96, 4.70, C.blue)
    for i, label in enumerate(["人口", "产业", "交通", "用地", "夜光", "POI"]):
        chip(s, label, 3.02 + i * 0.86, 4.66, 0.64, C.green, C.mint)
    text(s, "让模型从“看见地表变化”，进一步走向“理解变化背后的社会经济活动”。", 1.00, 5.32, 6.38, 0.22, 12, C.ink, True)
    for label, sub, x, color, fill in [
        ("基层治理", "更低门槛", 8.02, C.blue, C.pale_blue),
        ("中小企业", "按需调用", 9.64, C.green, C.mint),
        ("公众服务", "普惠应用", 11.26, C.purple, C.pale_purple),
    ]:
        rect(s, x, 4.86, 1.36, 0.70, fill, color)
        text(s, label, x + 0.08, 5.02, 1.20, 0.12, 11, color, True, PP_ALIGN.CENTER)
        text(s, sub, x + 0.08, 5.34, 1.20, 0.10, 7, C.body, True, PP_ALIGN.CENTER)
    statement(s, "长期愿景：让遥感能力像用地图一样简单，让地理智能成为社会经济运行的基础服务。", 6.58)

    return prs


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
