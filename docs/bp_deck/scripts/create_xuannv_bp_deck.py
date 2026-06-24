"""Generate the 22-page Xuannv investor BP deck.

This version follows docs/business_bp_final_slide_content.md page-by-page.
It keeps copy short and uses diagrams, image panels, and simple matrices instead
of dense paragraphs.
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "assets"
OUT_DIR = ROOT / "outputs"
OUT_PATH = OUT_DIR / "玄女科技BP_投资人版_v0.4_22页美化版.pptx"

COVER = ASSET_DIR / "geo_embedding_cover_pastel.png"
SPACE = ASSET_DIR / "space_data_to_intelligence.png"
GOV = ASSET_DIR / "government_change_review.png"
TEAM_WORK = ASSET_DIR / "remote_sensing_team_workflow.png"
LAB = ASSET_DIR / "university_research_lab.png"
OLD = ASSET_DIR / "old_bp_media"
OLD_EARTH = OLD / "image1.png"
OLD_OPTICAL = OLD / "image15.jpeg"
OLD_SAR = OLD / "image16.png"
OLD_EMBED = OLD / "image17.png"
OLD_LANDSLIDE = OLD / "image18.png"
OLD_STATION = OLD / "image28.png"
OLD_TEAM_1 = OLD / "image30.png"
OLD_TEAM_2 = OLD / "image31.jpeg"
OLD_TEAM_3 = OLD / "image32.png"
OLD_TEAM_4 = OLD / "image33.png"
VIS = Path(
    "/data/xuannv_embedding/outputs/downstream/visualizations/"
    "harbin_stage2_v1_fold0/patch_000021_visualization.png"
)


class C:
    bg = RGBColor(250, 248, 241)
    ink = RGBColor(34, 48, 53)
    muted = RGBColor(98, 111, 117)
    green = RGBColor(132, 181, 160)
    blue = RGBColor(136, 181, 205)
    lavender = RGBColor(180, 171, 210)
    sand = RGBColor(238, 229, 209)
    line = RGBColor(218, 226, 221)
    white = RGBColor(255, 255, 255)
    pale_green = RGBColor(242, 249, 245)
    pale_blue = RGBColor(241, 248, 251)
    pale_lav = RGBColor(247, 245, 251)
    dark_green = RGBColor(72, 121, 103)


FONT = "Microsoft YaHei"


def ensure_assets() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    if not VIS.exists():
        return
    image = Image.open(VIS).convert("RGB")
    panels = {
        "harbin_s2_rgb.png": (15, 120, 392, 497),
        "harbin_highres_optical.png": (459, 120, 836, 497),
        "harbin_s1_sar.png": (903, 120, 1280, 497),
        "harbin_landsat.png": (1347, 120, 1724, 497),
        "harbin_worldcover.png": (15, 528, 392, 905),
        "harbin_embedding_pca.png": (459, 528, 836, 905),
        "harbin_prediction.png": (903, 528, 1280, 905),
    }
    for name, box in panels.items():
        out = ASSET_DIR / name
        if not out.exists():
            image.crop(box).save(out)


def bg(slide, color=C.bg) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def text(
    slide,
    value: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 14,
    color=C.ink,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def title(slide, page: int, heading: str, sub: str = "") -> None:
    text(slide, f"{page:02d}", 0.72, 0.42, 0.45, 0.22, 8, C.dark_green, True)
    text(slide, heading, 1.18, 0.34, 10.8, 0.54, 23, C.ink, True)
    if sub:
        text(slide, sub, 1.2, 0.95, 10.7, 0.28, 10, C.muted)


def footer(slide, page: int) -> None:
    text(slide, "玄女科技商业计划书 v0.4｜22 页内容稿美化版", 0.72, 7.14, 6.2, 0.20, 7, C.muted)
    text(slide, f"{page:02d}", 12.05, 7.12, 0.5, 0.20, 8, C.muted, align=PP_ALIGN.RIGHT)


def rect(slide, x, y, w, h, fill=C.white, line=C.line, rounded=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(0.8)
    return s


def card(slide, head: str, body: str, x, y, w, h, accent=C.green, fill=C.white):
    rect(slide, x, y, w, h, fill)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    text(slide, head, x + 0.22, y + 0.16, w - 0.38, 0.25, 12, C.ink, True)
    text(slide, body, x + 0.22, y + 0.55, w - 0.38, h - 0.66, 9, C.muted)


def metric(slide, value: str, label: str, x, y, w, color=C.green):
    rect(slide, x, y, w, 0.92, RGBColor(253, 253, 250))
    text(slide, value, x + 0.15, y + 0.12, w - 0.3, 0.30, 18, color, True)
    text(slide, label, x + 0.15, y + 0.54, w - 0.3, 0.22, 8, C.muted)


def flow(slide, items: list[str], x, y, w, color=C.green):
    gap = 0.16
    bw = (w - gap * (len(items) - 1)) / len(items)
    for i, item in enumerate(items):
        bx = x + i * (bw + gap)
        rect(slide, bx, y, bw, 0.58, RGBColor(246, 250, 248), color)
        text(slide, item, bx + 0.06, y + 0.17, bw - 0.12, 0.18, 8, C.ink, True, PP_ALIGN.CENTER)
        if i < len(items) - 1:
            text(slide, "→", bx + bw + 0.02, y + 0.18, 0.1, 0.18, 9, C.muted, True)


def picture(slide, path: Path, x, y, w, h):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        rect(slide, x, y, w, h, C.pale_blue)
        text(slide, "素材待补", x, y + h / 2 - 0.12, w, 0.24, 11, C.muted, True, PP_ALIGN.CENTER)


def line(slide, x1, y1, x2, y2, color=C.line, width=1.2):
    shp = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shp.line.color.rgb = color
    shp.line.width = Pt(width)
    return shp


def bar_chart(slide, labels: list[str], values: list[float], x, y, w, h, max_value=None, colors=None, suffix=""):
    max_v = max_value or max(values)
    colors = colors or [C.green, C.blue, C.lavender, C.sand]
    line(slide, x, y + h, x + w, y + h, C.line, 1.0)
    gap = w / (len(values) * 2 + 1)
    bw = gap
    for i, (lab, val) in enumerate(zip(labels, values)):
        bx = x + gap + i * gap * 2
        bh = h * (val / max_v)
        rect(slide, bx, y + h - bh, bw, bh, colors[i % len(colors)], colors[i % len(colors)], rounded=False)
        text(slide, f"{val:g}{suffix}", bx - 0.18, y + h - bh - 0.25, bw + 0.36, 0.18, 8, C.ink, True, PP_ALIGN.CENTER)
        text(slide, lab, bx - 0.28, y + h + 0.10, bw + 0.56, 0.28, 7, C.muted, align=PP_ALIGN.CENTER)


def horizontal_bar(slide, labels: list[str], values: list[float], x, y, w, row_h=0.42, max_value=None, colors=None, suffix=""):
    max_v = max_value or max(values)
    colors = colors or [C.green, C.blue, C.lavender]
    for i, (lab, val) in enumerate(zip(labels, values)):
        yy = y + i * row_h
        text(slide, lab, x, yy + 0.05, 1.55, 0.18, 7, C.muted)
        rect(slide, x + 1.65, yy + 0.08, w, 0.14, RGBColor(238, 242, 240), C.line, rounded=False)
        rect(slide, x + 1.65, yy + 0.08, w * (val / max_v), 0.14, colors[i % len(colors)], colors[i % len(colors)], rounded=False)
        text(slide, f"{val:g}{suffix}", x + 1.75 + w * (val / max_v), yy + 0.03, 0.8, 0.18, 7, C.ink, True)


def pie_chart(slide, parts: list[tuple[str, float, RGBColor]], x, y, size):
    start = 0
    total = sum(v for _, v, _ in parts)
    for lab, val, color in parts:
        end = start + int(360 * val / total)
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PIE, Inches(x), Inches(y), Inches(size), Inches(size))
        shp.adjustments[0] = start
        shp.adjustments[1] = end
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.color.rgb = C.bg
        start = end
    rect(slide, x + size * 0.32, y + size * 0.32, size * 0.36, size * 0.36, C.bg, C.bg)
    for i, (lab, val, color) in enumerate(parts):
        yy = y + 0.18 + i * 0.38
        rect(slide, x + size + 0.35, yy, 0.16, 0.16, color, color, rounded=False)
        text(slide, f"{lab} {val:g}%", x + size + 0.58, yy - 0.03, 2.1, 0.18, 8, C.muted)


def bubble(slide, label: str, value: str, x, y, size, color):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.fill.transparency = 12
    shp.line.color.rgb = color
    text(slide, value, x, y + size * 0.32, size, 0.22, 13, C.ink, True, PP_ALIGN.CENTER)
    text(slide, label, x + 0.08, y + size * 0.56, size - 0.16, 0.25, 7, C.muted, align=PP_ALIGN.CENTER)


def layered_stack(slide, layers: list[tuple[str, str, RGBColor]], x, y, w):
    for i, (head, body, color) in enumerate(layers):
        yy = y + i * 1.05
        rect(slide, x + i * 0.28, yy, w - i * 0.56, 0.82, RGBColor(253, 253, 250), color)
        text(slide, head, x + i * 0.28 + 0.18, yy + 0.12, w - i * 0.56 - 0.36, 0.18, 11, C.ink, True, PP_ALIGN.CENTER)
        text(slide, body, x + i * 0.28 + 0.18, yy + 0.38, w - i * 0.56 - 0.36, 0.18, 7, C.muted, align=PP_ALIGN.CENTER)


def cost_curve(slide, x, y, w, h):
    line(slide, x, y + h, x + w, y + h, C.line, 1)
    line(slide, x, y, x, y + h, C.line, 1)
    # Traditional curve.
    pts1 = [(x + 0.3, y + h - 0.45), (x + w * 0.36, y + h - 1.35), (x + w * 0.68, y + h - 2.25), (x + w - 0.25, y + 0.35)]
    pts2 = [(x + 0.3, y + h - 0.7), (x + w * 0.36, y + h - 1.0), (x + w * 0.68, y + h - 1.25), (x + w - 0.25, y + h - 1.45)]
    for pts, color, label in [(pts1, C.blue, "传统：任务越多成本线性上升"), (pts2, C.green, "玄女：底座复用后边际成本下降")]:
        for a, b in zip(pts, pts[1:]):
            line(slide, a[0], a[1], b[0], b[1], color, 2.4)
        text(slide, label, pts[-1][0] - 2.65, pts[-1][1] - 0.28, 2.55, 0.20, 8, color, True)
    text(slide, "任务数", x + w - 0.55, y + h + 0.14, 0.5, 0.18, 7, C.muted)
    text(slide, "总成本", x - 0.15, y - 0.16, 0.6, 0.18, 7, C.muted)


def mini_table(slide, rows: list[list[str]], x, y, col_w: list[float], row_h=0.48, size=7):
    total = sum(col_w)
    for r, row in enumerate(rows):
        fill = C.pale_green if r == 0 else C.white
        rect(slide, x, y + r * row_h, total, row_h - 0.03, fill, C.line, rounded=False)
        cx = x
        for c, cell in enumerate(row):
            text(slide, cell, cx + 0.06, y + r * row_h + 0.10, col_w[c] - 0.1, row_h - 0.16, size, C.ink if r == 0 else C.muted, r == 0 or c == 0)
            cx += col_w[c]


def three_points(slide, points: list[tuple[str, str]], y=4.85):
    colors = [C.green, C.blue, C.lavender]
    for i, (h, b) in enumerate(points):
        card(slide, h, b, 0.82 + i * 4.05, y, 3.55, 1.24, colors[i])


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank); bg(s)
    picture(s, COVER, 5.15, 0, 8.2, 7.5)
    text(s, "玄女科技", 0.82, 1.22, 4.8, 0.48, 24, C.ink, True)
    text(s, "用地理嵌入赋能遥感智能应用", 0.82, 1.88, 5.8, 0.72, 27, C.dark_green, True)
    text(s, "一次生成地球表征，多任务复用。\n让中国拥有自己的地理智能底座。", 0.86, 3.02, 5.1, 0.75, 15, C.muted)
    text(s, "商业计划书｜融资需求 5000 万元", 0.86, 6.45, 4.8, 0.25, 10, C.muted)

    # 2
    s = prs.slides.add_slide(blank); bg(s); title(s, 2, "美国已经证明：地理智能可以长成大公司", "不同标杆证明不同环节：数据、监测、平台、地理嵌入趋势。")
    picture(s, OLD_EARTH, 8.45, 1.18, 3.65, 2.05)
    bubble(s, "Maxar\n收购估值", "64 亿美元", 0.95, 1.55, 1.55, C.green)
    bubble(s, "Planet\nFY2025 收入", "2.444 亿", 2.85, 1.35, 1.85, C.blue)
    bubble(s, "BlackSky\n2024 收入", "1.021 亿", 5.05, 1.72, 1.35, C.lavender)
    bar_chart(s, ["数据", "监测", "平台", "嵌入"], [4, 3, 5, 2], 0.95, 4.28, 5.8, 1.05, max_value=5)
    mini_table(s, [["标杆", "证明什么"], ["Maxar / Planet / BlackSky", "影像与实时监测存在持续需求"], ["Esri / CARTO", "地理平台可以进入政企工作流"], ["Google 卫星嵌入", "表征层正在成为下一代入口"]], 7.25, 3.78, [2.35, 3.35], 0.52, 8)
    text(s, "玄女借鉴的是“AI 可用地理表征层”，不是复制任何一家美国公司的资产结构。", 1.0, 6.2, 11.2, 0.35, 15, C.dark_green, True, PP_ALIGN.CENTER)
    footer(s, 2)

    # 3
    s = prs.slides.add_slide(blank); bg(s); title(s, 3, "行业正在从“卖影像”走向“卖可调用的地球智能”")
    picture(s, SPACE, 7.15, 1.28, 5.25, 3.7)
    stages = [("影像时代", "卖原始影像\n底图 / 授权"), ("监测时代", "卖变化监测\n目标识别 / 预警"), ("平台时代", "卖 GIS / 云平台\n接口 / 工作流"), ("嵌入时代", "卖可复用表征\n分类 / 检测 / 问答")]
    for i, (h, b) in enumerate(stages):
        x = 0.85 + i * 1.55
        rect(s, x, 2.0, 1.22, 1.22, [C.sand, C.pale_blue, C.pale_green, C.pale_lav][i], [C.green, C.blue, C.green, C.lavender][i])
        text(s, str(i + 1), x + 0.43, 2.12, 0.32, 0.22, 12, [C.green, C.blue, C.green, C.lavender][i], True, PP_ALIGN.CENTER)
        text(s, h, x - 0.05, 3.42, 1.35, 0.22, 9, C.ink, True, PP_ALIGN.CENTER)
        text(s, b, x - 0.15, 3.78, 1.55, 0.45, 7, C.muted, align=PP_ALIGN.CENTER)
        if i < 3:
            line(s, x + 1.22, 2.62, x + 1.55, 2.62, C.line, 1.6)
    card(s, "玄女的位置", "不再只交付影像或单点分析，而是切入更靠近应用价值的 AI 可用地理表征层。", 0.95, 5.28, 5.75, 0.95, C.green)
    footer(s, 3)

    # 4
    s = prs.slides.add_slide(blank); bg(s); title(s, 4, "中国也到了这个窗口：商业航天解决“看见”，地理智能解决“理解”", "市场数字为广义产业和上游供给口径，不等同于玄女直接可获得收入。")
    bubble(s, "中国地理信息产业\n广义基础", "8501 亿元", 0.95, 1.62, 1.9, C.green)
    bubble(s, "EO 小卫星\n上游供给", "55.2 亿美元", 3.45, 1.52, 2.05, C.blue)
    bubble(s, "地理空间分析\n广义大盘", "2340 亿美元", 6.1, 1.35, 2.35, C.lavender)
    bar_chart(s, ["2025", "2030"], [26.4, 55.2], 9.35, 1.65, 2.55, 1.5, suffix="")
    text(s, "EO 小卫星市场增长", 9.25, 3.48, 2.8, 0.20, 8, C.muted, align=PP_ALIGN.CENTER)
    three_points(s, [("数据供给", "商业航天、高分光学、SAR、无人机与北斗持续增长。"), ("场景密集", "自然资源、应急、农业、水利、能源、重大工程都要监测。"), ("支付基础", "已有政企采购和项目交付基础，需要更高复用效率。")], 4.25)
    text(s, "中国的机会不是再多一批影像，而是把影像变成可复用的地理智能。", 1.0, 6.25, 11.0, 0.35, 17, C.dark_green, True, PP_ALIGN.CENTER)
    footer(s, 4)

    # 5
    s = prs.slides.add_slide(blank); bg(s); title(s, 5, "今天的遥感公司，主要还是靠项目、人力和数据交付赚钱")
    rows = [["收入方式", "受限点"], ["数据销售", "价值停在数据层"], ["项目交付", "项目越多人越多"], ["平台授权", "智能任务仍需定制"], ["专业服务", "毛利被人力吞掉"], ["订阅接口", "行业仍在早期"]]
    mini_table(s, rows, 0.95, 1.45, [2.25, 3.25], 0.55, 8)
    pie_chart(s, [("数据销售", 18, C.green), ("项目交付", 42, C.blue), ("平台授权", 18, C.lavender), ("专业服务", 17, C.sand), ("订阅接口", 5, RGBColor(190, 212, 203))], 7.2, 1.55, 2.2)
    text(s, "示意：传统收入仍以项目和服务为主", 7.1, 4.05, 4.4, 0.2, 9, C.muted, align=PP_ALIGN.CENTER)
    text(s, "行业不是没有需求，而是需求越多，传统交付方式越重。", 1.0, 6.15, 11.1, 0.38, 19, C.dark_green, True, PP_ALIGN.CENTER)
    footer(s, 5)

    # 6
    s = prs.slides.add_slide(blank); bg(s); title(s, 6, "真正的瓶颈不是数据少，而是每个任务都重做一遍")
    flow(s, ["找数据", "预处理", "标注", "训练", "核查", "交付"], 0.95, 1.55, 11.2, C.blue)
    bar_chart(s, ["标注效率", "年省成本", "三年效益"], [80, 200, 5000], 1.2, 3.35, 4.2, 1.65, max_value=5000, suffix="")
    metric(s, "80%", "AI 辅助标注效率提升", 6.45, 3.1, 2.75, C.green)
    metric(s, "200 万元+", "年均节省人力成本", 9.55, 3.1, 2.75, C.blue)
    text(s, "该案例证明标注环节有明确降本空间；玄女要把这种效率提升从单点工具扩展到多任务复用。", 1.0, 5.5, 11.1, 0.55, 18, C.ink, True, PP_ALIGN.CENTER)
    footer(s, 6)

    # 7
    s = prs.slides.add_slide(blank); bg(s); title(s, 7, "真实痛点 1：遥感项目公司不是缺项目，而是项目越多越像人力外包")
    picture(s, TEAM_WORK, 0.78, 1.35, 5.6, 3.95)
    card(s, "今天怎么做", "施工、农用地、违建、裸地等项目各开一条流程。", 6.78, 1.55, 5.05, 0.96, C.blue)
    card(s, "真正痛点", "多项目并行、人手不够；标注和质检成本高；跨区域迁移不稳定。", 6.78, 2.8, 5.05, 1.06, C.lavender)
    card(s, "玄女价值", "区域地理嵌入底座 + 少样本适配 + 变化候选和图斑预筛。", 6.78, 4.15, 5.05, 1.06, C.green)
    text(s, "玄女卖的不是一张图，而是把项目制遥感交付的重复劳动压缩成可复用的模型资产。", 0.95, 6.12, 11.5, 0.35, 16, C.dark_green, True, PP_ALIGN.CENTER)
    footer(s, 7)

    # 8
    s = prs.slides.add_slide(blank); bg(s); title(s, 8, "真实痛点 2：政府不是缺卫星图，而是缺能进入流程的变化证据链")
    picture(s, GOV, 6.85, 1.3, 5.65, 4.05)
    flow(s, ["疑似图斑", "外业核查", "举证判定", "整改复核", "归档闭环"], 0.9, 1.75, 5.55, C.green)
    card(s, "痛点", "图斑多、专项任务交叉、时间集中、遥感结果不能直接执法。", 0.95, 3.15, 5.35, 0.98, C.blue)
    card(s, "玄女", "图斑优先级、变化类型理解、历史过程证据链、冲突点提示。", 0.95, 4.45, 5.35, 0.98, C.lavender)
    footer(s, 8)

    # 9
    s = prs.slides.add_slide(blank); bg(s); title(s, 9, "真实痛点 3：普通地块用户不想学遥感，只想知道自己的地有没有异常")
    metric(s, "27%", "2023 年美国农场/牧场使用精准农业实践", 0.95, 1.55, 3.6, C.green)
    card(s, "今天", "人工巡查、拍照反馈、一次性报告；不会处理专业遥感图层。", 0.95, 3.0, 3.6, 1.15, C.blue)
    card(s, "玄女", "圈定地块，监测硬化、堆土、棚房、道路、积水、作物异常。", 4.95, 3.0, 3.6, 1.15, C.green)
    card(s, "定位", "不是近期主收入，但说明地理嵌入能让遥感能力普惠化。", 8.95, 3.0, 3.6, 1.15, C.lavender)
    flow(s, ["圈地块", "自动监测", "异常提醒", "证据包"], 2.0, 5.3, 9.3, C.green)
    footer(s, 9)

    # 10
    s = prs.slides.add_slide(blank); bg(s); title(s, 10, "真实痛点 4：高校和科研团队有想法，但数据工程吃掉研究周期")
    picture(s, LAB, 0.82, 1.28, 5.65, 4.0)
    card(s, "今天", "下载、配准、裁切、标注、转换格式、跑环境；样本制作本身就是大工程。", 6.85, 1.5, 5.0, 1.1, C.blue)
    card(s, "痛点", "数据准备比算法研究更耗时；小团队缺算力；跨区域复现难。", 6.85, 2.92, 5.0, 1.1, C.lavender)
    card(s, "玄女", "现成地理嵌入、预标注、变化候选、难例检索、少样本微调。", 6.85, 4.34, 5.0, 1.1, C.green)
    footer(s, 10)

    # 11
    s = prs.slides.add_slide(blank); bg(s); title(s, 11, "从四类用户出发，我们的核心问题是什么？")
    # Four-quadrant pain map.
    rect(s, 1.15, 1.55, 10.8, 3.5, C.white, C.line)
    line(s, 6.55, 1.55, 6.55, 5.05, C.line, 1)
    line(s, 1.15, 3.3, 11.95, 3.3, C.line, 1)
    text(s, "B 端遥感公司\n项目越多，越依赖人力交付", 1.45, 2.12, 4.6, 0.45, 15, C.blue, True, PP_ALIGN.CENTER)
    text(s, "政府部门\n图斑越多，越需要证据链", 6.95, 2.12, 4.6, 0.45, 15, C.green, True, PP_ALIGN.CENTER)
    text(s, "地块用户\n有地理问题，但不会用遥感", 1.45, 3.82, 4.6, 0.45, 15, C.lavender, True, PP_ALIGN.CENTER)
    text(s, "科研团队\n想法很多，但数据工程太重", 6.95, 3.82, 4.6, 0.45, 15, C.dark_green, True, PP_ALIGN.CENTER)
    text(s, "如何把爆发式增长的地球观测数据转化为可复用的地理智能底座？", 1.0, 5.75, 11.1, 0.44, 22, C.ink, True, PP_ALIGN.CENTER)
    footer(s, 11)

    # 12
    s = prs.slides.add_slide(blank); bg(s); title(s, 12, "玄女的答案：先生成地理嵌入，再服务下游任务")
    flow(s, ["多源观测", "地理嵌入", "任务组件", "变化候选", "业务交付"], 0.95, 1.42, 11.25, C.green)
    # Hub-and-spoke task map.
    center_x, center_y = 5.7, 3.55
    bubble(s, "统一表征", "地理嵌入", center_x, center_y, 1.55, C.green)
    spokes = [("施工工地", 2.0, 3.0, C.blue), ("建筑变化", 3.55, 4.85, C.lavender), ("农用地变化", 7.85, 4.85, C.green), ("灾害线索", 9.4, 3.0, C.blue), ("缺失补齐", 5.75, 5.45, C.sand)]
    for lab, x, y, col in spokes:
        line(s, center_x + 0.78, center_y + 0.78, x + 0.55, y + 0.55, C.line, 1.0)
        bubble(s, lab, "", x, y, 1.1, col)
    text(s, "先把地球变成机器能复用的表示，再让所有任务调用这个表示。", 1.0, 6.45, 11.2, 0.26, 16, C.dark_green, True, PP_ALIGN.CENTER)
    footer(s, 12)

    # 13
    s = prs.slides.add_slide(blank); bg(s); title(s, 13, "地理嵌入是遥感行业的 token 层")
    card(s, "文本智能", "文本 → token → 理解 / 检索 / 生成", 1.2, 1.7, 4.6, 1.25, C.blue)
    card(s, "地球智能", "地球观测 → 地理嵌入 → 变化检测 / 分类 / 预测 / 问答", 7.15, 1.7, 4.6, 1.25, C.green)
    text(s, "GPT 先把文本转成 token，才能进行理解、检索、生成；玄女先把地球观测转成地理嵌入，才能进行变化检测、分类、预测和问答。", 1.15, 3.85, 11.1, 0.9, 24, C.ink, True, PP_ALIGN.CENTER)
    text(s, "商业航天解决“看见地球”，地理嵌入解决“理解地球”。", 1.0, 5.8, 11.2, 0.36, 18, C.dark_green, True, PP_ALIGN.CENTER)
    footer(s, 13)

    # 14
    s = prs.slides.add_slide(blank); bg(s); title(s, 14, "一次嵌入，多任务复用，改变遥感交付的成本曲线")
    cost_curve(s, 1.1, 1.55, 5.25, 3.65)
    card(s, "传统模式", "总成本 ≈ 任务数 ×（数据处理 + 标注 + 建模 + 推理 + 核查 + 报告）", 7.05, 1.68, 4.9, 1.0, C.blue)
    card(s, "玄女模式", "总成本 ≈ 区域底座成本 + 任务数 ×（少量样本 + 轻量适配 + 核查报告）", 7.05, 2.95, 4.9, 1.0, C.green)
    horizontal_bar(s, ["时间", "存储", "算力", "人力"], [75, 60, 55, 70], 7.15, 4.4, 2.7, 0.42, max_value=100, suffix="%")
    text(s, "第一个任务验证价值，第二个任务开始复用；任务越多，底座越值钱。", 1.0, 6.05, 11.2, 0.36, 19, C.dark_green, True, PP_ALIGN.CENTER)
    footer(s, 14)

    # 15
    s = prs.slides.add_slide(blank); bg(s); title(s, 15, "产品形态：不是一个模型，而是城市地理智能底座")
    layered_stack(s, [("业务交付层", "图斑 / 证据链 / 报告 / 接口 / 驾驶舱", C.lavender), ("任务组件库", "施工 / 建筑 / 农用地 / 灾害 / 生态", C.blue), ("地理嵌入底座", "按区域、按月份持续生成可复用表征", C.green)], 1.25, 1.55, 6.0)
    rows = [["客户", "产品"], ["政府 / 新区", "年度变化监测"], ["遥感公司", "任务库授权 / 接口"], ["学校科研", "嵌入数据集 / 开发生态"], ["地块用户", "异常提醒 / 证据报告"]]
    mini_table(s, rows, 8.0, 1.65, [1.6, 2.85], 0.52, 8)
    picture(s, OLD_STATION, 8.2, 4.55, 3.9, 1.55)
    footer(s, 15)

    # 16
    s = prs.slides.add_slide(blank); bg(s); title(s, 16, "样板 1：哈尔滨新区，用于验证“一区多任务复用”")
    labels = [("S2 影像", "harbin_s2_rgb.png"), ("高分光学", "harbin_highres_optical.png"), ("地理嵌入", "harbin_embedding_pca.png"), ("预测概率", "harbin_prediction.png")]
    for i, (lab, name) in enumerate(labels):
        x = 0.95 + i * 2.05
        text(s, lab, x, 1.45, 1.7, 0.18, 7, C.muted, True, PP_ALIGN.CENTER)
        picture(s, ASSET_DIR / name, x, 1.72, 1.75, 1.75)
    card(s, "验证什么", "一个区域生成月度地理嵌入；同一套表征支撑施工、建筑、农用地、垃圾/裸地等任务。", 9.35, 1.55, 2.85, 1.72, C.green)
    flow(s, ["区域底座", "施工", "建筑变化", "农用地", "异常变化"], 1.1, 4.65, 10.8, C.green)
    footer(s, 16)

    # 17
    s = prs.slides.add_slide(blank); bg(s); title(s, 17, "样板 2：海淀区，用于验证复杂城市纹理和高分模态融合")
    picture(s, ASSET_DIR / "harbin_highres_optical.png", 0.95, 1.55, 3.1, 3.1)
    picture(s, ASSET_DIR / "harbin_s1_sar.png", 4.45, 1.55, 3.1, 3.1)
    card(s, "验证任务", "高密度城市空间、小目标、复杂边界；高分光学 / 高分 SAR 与低分时序融合。", 8.1, 1.62, 4.05, 1.35, C.blue)
    card(s, "输出证据", "高分影像、标注样本、模型输出与人工标注对比、模态增强前后差异。", 8.1, 3.35, 4.05, 1.35, C.green)
    text(s, "海淀页当前用已有高分样例占位；后续可替换为真实海淀最佳 patch。", 1.0, 5.82, 11.2, 0.35, 13, C.muted, align=PP_ALIGN.CENTER)
    footer(s, 17)

    # 18
    s = prs.slides.add_slide(blank); bg(s); title(s, 18, "样板 3：雅江，用于验证极端地形、灾害场景和缺失模态补齐")
    picture(s, ASSET_DIR / "harbin_landsat.png", 0.95, 1.55, 3.05, 3.05)
    picture(s, ASSET_DIR / "harbin_s1_sar.png", 4.35, 1.55, 3.05, 3.05)
    card(s, "复杂性", "山地峡谷、河谷、灾害风险、施工扰动；光学云雾遮挡和模态缺失。", 8.0, 1.55, 4.25, 1.2, C.lavender)
    card(s, "验证任务", "滑坡、崩塌、堰塞湖、施工扰动候选；光学缺失时用 SAR / 时序补齐。", 8.0, 3.05, 4.25, 1.2, C.green)
    card(s, "需要补充", "真实雅江数据目录、云遮挡样例、缺失模态补齐结果、人工核查命中率。", 8.0, 4.55, 4.25, 1.2, C.blue)
    footer(s, 18)

    # 19
    s = prs.slides.add_slide(blank); bg(s); title(s, 19, "能力展示：嵌入、补齐、时序、多任务迁移")
    panels = [("原始影像", "harbin_highres_optical.png"), ("地理嵌入", "harbin_embedding_pca.png"), ("预测概率", "harbin_prediction.png"), ("土地覆盖", "harbin_worldcover.png")]
    for i, (lab, name) in enumerate(panels):
        x = 0.9 + i * 2.25
        text(s, lab, x, 1.45, 1.85, 0.18, 7, C.muted, True, PP_ALIGN.CENTER)
        picture(s, ASSET_DIR / name, x, 1.75, 1.85, 1.85)
    card(s, "四类能力", "地理嵌入｜缺失模态补齐｜多源时序理解｜下游迁移", 1.0, 4.55, 5.4, 1.1, C.green)
    card(s, "核心判断", "玄女的核心不是某个任务准确率，而是一个表征可以服务多少任务。", 6.95, 4.55, 5.1, 1.1, C.blue)
    footer(s, 19)

    # 20
    s = prs.slides.add_slide(blank); bg(s); title(s, 20, "为什么我们能做：多源时空地理嵌入模型架构")
    picture(s, OLD_OPTICAL, 0.85, 1.38, 2.05, 1.15)
    picture(s, OLD_SAR, 0.85, 2.75, 2.05, 1.05)
    picture(s, OLD_EMBED, 0.85, 4.0, 2.05, 1.05)
    flow(s, ["多源输入", "时空对齐", "编码器", "地理嵌入", "任务头"], 3.35, 1.9, 8.5, C.green)
    layered_stack(s, [("输入", "S2 / S1 / Landsat / 高分 / SAR / DEM / 气象", C.blue), ("表征", "跨时间、跨传感器的稳定地物表示", C.green), ("输出", "补齐 / 检测 / 分类 / 分割 / 问答", C.lavender)], 3.65, 3.35, 7.9)
    text(s, "我们不是把多源影像简单拼接，而是让模型学习同一地点在不同时间、不同传感器下的稳定表征。", 1.0, 5.95, 11.2, 0.4, 17, C.dark_green, True, PP_ALIGN.CENTER)
    footer(s, 20)

    # 21
    s = prs.slides.add_slide(blank); bg(s); title(s, 21, "商业模式：先用标杆项目拿信任，再用底座和任务库拿复用收入")
    # Revenue flywheel.
    bubble(s, "政府合同", "压舱石", 1.2, 1.75, 1.55, C.green)
    bubble(s, "B 端授权", "现金流", 4.05, 1.75, 1.55, C.blue)
    bubble(s, "接口生态", "放大器", 6.9, 1.75, 1.55, C.lavender)
    bubble(s, "数据反馈", "护城河", 9.75, 1.75, 1.55, C.sand)
    line(s, 2.75, 2.52, 4.05, 2.52, C.line, 1.5)
    line(s, 5.6, 2.52, 6.9, 2.52, C.line, 1.5)
    line(s, 8.45, 2.52, 9.75, 2.52, C.line, 1.5)
    rows = [["阶段", "收入方式"], ["近期", "标杆项目 / 年度监测"], ["中期", "城市底座订阅 / 任务库授权"], ["长期", "接口调用 / 数据分成 / 应用市场"]]
    mini_table(s, rows, 2.25, 4.35, [2.0, 6.8], 0.5, 9)
    text(s, "政府合同是压舱石，B 端授权是现金流，接口和生态是放大器。", 1.0, 6.35, 11.2, 0.32, 16, C.dark_green, True, PP_ALIGN.CENTER)
    footer(s, 21)

    # 22
    s = prs.slides.add_slide(blank); bg(s); title(s, 22, "融资 5000 万：把技术样板推成商业样板")
    metric(s, "5000 万", "融资金额", 0.95, 1.38, 2.7, C.green)
    metric(s, "18 个月", "样板和商业化周期", 3.95, 1.38, 2.7, C.blue)
    metric(s, "5-8 个", "可复用下游任务", 6.95, 1.38, 2.7, C.lavender)
    pie_chart(s, [("模型数据", 35, C.green), ("国产算力", 20, C.blue), ("任务库", 20, C.lavender), ("样板交付", 15, C.sand), ("商务团队", 10, RGBColor(190, 212, 203))], 10.05, 1.25, 1.55)
    rows = [["里程碑", "目标"], ["城市级样板", "哈尔滨新区多任务连续监测"], ["补充场景", "海淀复杂城市 / 雅江复杂地形"], ["商业结果", "1-2 个年度服务或 B 端授权意向 / 订单"], ["交付能力", "接口、任务库、报告流程可演示"]]
    mini_table(s, rows, 1.0, 3.25, [3.0, 6.2], 0.5, 8)
    text(s, "这轮融资不是为了做一组模型指标，而是为了证明地理嵌入底座可以规模化交付。", 1.0, 6.25, 11.2, 0.35, 17, C.dark_green, True, PP_ALIGN.CENTER)
    footer(s, 22)

    return prs


def main() -> None:
    ensure_assets()
    prs = build()
    prs.save(OUT_PATH)
    print(OUT_PATH)


if __name__ == "__main__":
    main()
